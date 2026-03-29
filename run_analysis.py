"""
Chicago Traffic Crash Analysis (2016–2023)
Complete analysis script — generates charts, maps, stats, and PowerPoint.
"""
import warnings
warnings.filterwarnings('ignore')

import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as ticker
import seaborn as sns
import scipy
import scipy.stats as stats
import statsmodels.api as sm
import statsmodels.formula.api as smf
from statsmodels.tsa.seasonal import seasonal_decompose
from sklearn.cluster import DBSCAN
import folium
from folium.plugins import HeatMap
import sqlite3
import calendar
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Configure plotting
sns.set_theme(style='whitegrid')
plt.rcParams.update({
    'figure.dpi': 250,
    'savefig.dpi': 250,
    'figure.facecolor': 'white',
    'savefig.facecolor': 'white',
    'font.family': 'sans-serif',
    'font.size': 11,
    'axes.titlesize': 14,
    'axes.labelsize': 12,
})

OUTPUT_DIR = 'output'
os.makedirs(OUTPUT_DIR, exist_ok=True)

print("All libraries loaded successfully.")
print(f"pandas {pd.__version__}, numpy {np.__version__}, scipy {scipy.__version__}")
print(f"matplotlib {matplotlib.__version__}, seaborn {sns.__version__}")
print(f"statsmodels {sm.__version__}")

# ============================================================
# PHASE 1: Data Loading, Cleaning & Feature Engineering
# ============================================================
print("\n" + "=" * 70)
print("PHASE 1: DATA LOADING, CLEANING & FEATURE ENGINEERING")
print("=" * 70)

# 1.1 Load CSV
df_raw = pd.read_csv('Traffic_Crashes.csv', parse_dates=['CRASH_DATE'], low_memory=False)
print(f"Raw dataset: {df_raw.shape[0]:,} rows, {df_raw.shape[1]} columns")
print(f"Date range: {df_raw['CRASH_DATE'].min()} to {df_raw['CRASH_DATE'].max()}")

# 1.2 Data Quality
null_counts = df_raw.isnull().sum()
null_pct = (null_counts / len(df_raw) * 100).round(1)
null_summary = pd.DataFrame({'Missing': null_counts, '% Missing': null_pct})
null_summary = null_summary[null_summary['Missing'] > 0].sort_values('% Missing', ascending=False)
print("\nColumns with missing values:")
print(null_summary.to_string())

print(f"\nCRASH_HOUR range: {df_raw['CRASH_HOUR'].min()} - {df_raw['CRASH_HOUR'].max()}")
print(f"CRASH_DAY_OF_WEEK range: {df_raw['CRASH_DAY_OF_WEEK'].min()} - {df_raw['CRASH_DAY_OF_WEEK'].max()}")
print(f"CRASH_MONTH range: {df_raw['CRASH_MONTH'].min()} - {df_raw['CRASH_MONTH'].max()}")

print(f"\nHIT_AND_RUN_I value counts:")
print(df_raw['HIT_AND_RUN_I'].value_counts(dropna=False).to_string())

# 1.3 Filter & Clean
df = df_raw.copy()
df['CRASH_YEAR'] = df['CRASH_DATE'].dt.year

print(f"\nYear distribution before filter:")
print(df['CRASH_YEAR'].value_counts().sort_index().to_string())

df = df[df['CRASH_YEAR'].between(2016, 2023)].copy()
print(f"\nAfter filtering to 2016-2023: {len(df):,} rows")

df = df.drop(columns=['DOORING_I', 'LANE_CNT'], errors='ignore')
print("Dropped DOORING_I (99.7% blank) and LANE_CNT (74.7% missing)")

df['HIT_AND_RUN_I'] = df['HIT_AND_RUN_I'].fillna('N')
print(f"\nHIT_AND_RUN_I after cleaning:")
print(df['HIT_AND_RUN_I'].value_counts().to_string())

before = len(df)
df = df.dropna(subset=['CRASH_HOUR', 'CRASH_DAY_OF_WEEK'])
print(f"\nDropped {before - len(df)} rows with null CRASH_HOUR/CRASH_DAY_OF_WEEK")
print(f"Final cleaned dataset: {len(df):,} rows")

# 1.4 Feature Engineering
wet_weather = ['RAIN', 'SNOW', 'FREEZING RAIN/DRIZZLE', 'SLEET/HAIL', 'BLOWING SNOW']
wet_surface = ['WET', 'SNOW OR SLUSH', 'ICE']

df['WEATHER_KNOWN'] = df['WEATHER_CONDITION'].notna() & ~df['WEATHER_CONDITION'].isin(['UNKNOWN'])
df['IS_WET_WEATHER'] = df['WEATHER_CONDITION'].isin(wet_weather).astype(int)
df['IS_WET_SURFACE'] = df['ROADWAY_SURFACE_COND'].isin(wet_surface).astype(int)
df['HAS_INJURY'] = (df['INJURIES_TOTAL'] > 0).astype(int)
df['HAS_FATAL'] = (df['INJURIES_FATAL'] > 0).astype(int)
severe_injuries = ['INCAPACITATING INJURY', 'FATAL']
df['IS_SEVERE'] = df['MOST_SEVERE_INJURY'].isin(severe_injuries).astype(int)
df['IS_WEEKEND'] = df['CRASH_DAY_OF_WEEK'].isin([1, 7]).astype(int)

def assign_time_period(hour):
    if 6 <= hour <= 9: return 'Morning Rush'
    elif 10 <= hour <= 15: return 'Midday'
    elif 16 <= hour <= 19: return 'Evening Rush'
    elif 20 <= hour <= 23: return 'Night'
    else: return 'Late Night'

df['TIME_PERIOD'] = df['CRASH_HOUR'].apply(assign_time_period)

day_map = {1: 'Sunday', 2: 'Monday', 3: 'Tuesday', 4: 'Wednesday',
           5: 'Thursday', 6: 'Friday', 7: 'Saturday'}
df['DAY_NAME'] = df['CRASH_DAY_OF_WEEK'].map(day_map)
df['MONTH_NAME'] = df['CRASH_MONTH'].apply(lambda m: calendar.month_abbr[int(m)])

def speed_cat(s):
    if pd.isna(s): return 'Unknown'
    if s <= 20: return '0-20 mph'
    if s <= 30: return '21-30 mph'
    if s <= 40: return '31-40 mph'
    return '45+ mph'  # Chicago speed limits are in 5 mph increments, so 41-44 is effectively empty

df['SPEED_CATEGORY'] = df['POSTED_SPEED_LIMIT'].apply(speed_cat)

damage_map = {'$500 OR LESS': 1, '$501 - $1,500': 2, 'OVER $1,500': 3}
df['DAMAGE_LEVEL'] = df['DAMAGE'].map(damage_map)
df['IS_DARK'] = df['LIGHTING_CONDITION'].str.contains('DARK', na=False).astype(int)

print("\nFeature engineering complete.")
eng_cols = ['IS_WET_WEATHER', 'IS_WET_SURFACE', 'HAS_INJURY', 'HAS_FATAL',
            'IS_SEVERE', 'IS_WEEKEND', 'IS_DARK']
for col in eng_cols:
    print(f"  {col}: {df[col].sum():,} ({df[col].mean()*100:.1f}%)")

# 1.5 Save Cleaned Data
try:
    df.to_parquet('crashes_clean.parquet', index=False)
    print("\nSaved to crashes_clean.parquet")
except Exception as e:
    print(f"\nParquet save failed ({e}), falling back to CSV")
    df.to_csv('crashes_clean.csv', index=False)

# ============================================================
# PHASE 2: SQL-Based Exploration
# ============================================================
print("\n" + "=" * 70)
print("PHASE 2: SQL-BASED EXPLORATION (SQLite)")
print("=" * 70)

conn = sqlite3.connect('crashes.db')
df.to_sql('crashes', conn, if_exists='replace', index=False)
print(f"Loaded {len(df):,} rows into SQLite database 'crashes.db'")

# Query 1: Crashes by Year
query1 = """
SELECT CRASH_YEAR, COUNT(*) as total_crashes,
       SUM(INJURIES_TOTAL) as total_injuries,
       SUM(INJURIES_FATAL) as total_fatalities,
       ROUND(AVG(HAS_INJURY)*100, 1) as injury_rate_pct,
       ROUND(AVG(HAS_FATAL)*100, 3) as fatality_rate_pct
FROM crashes GROUP BY CRASH_YEAR ORDER BY CRASH_YEAR
"""
result1 = pd.read_sql_query(query1, conn)
print("\nQuery 1 — Crashes by Year:")
print(result1.to_string(index=False))

# Query 2: Day × Hour Cross-Tab
query2 = """
SELECT DAY_NAME, CRASH_DAY_OF_WEEK, CRASH_HOUR, COUNT(*) as crash_count
FROM crashes GROUP BY CRASH_DAY_OF_WEEK, CRASH_HOUR
ORDER BY CRASH_DAY_OF_WEEK, CRASH_HOUR
"""
result2 = pd.read_sql_query(query2, conn)
heatmap_data = result2.pivot(index='DAY_NAME', columns='CRASH_HOUR', values='crash_count')
day_order = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday', 'Sunday']
heatmap_data = heatmap_data.reindex(day_order)
print("\nQuery 2 — Day x Hour cross-tab generated for heatmap.")

# Query 3: Weather Impact
query3 = """
SELECT WEATHER_CONDITION, COUNT(*) as total_crashes,
       SUM(HAS_INJURY) as injury_crashes,
       ROUND(AVG(HAS_INJURY)*100, 1) as injury_rate_pct,
       SUM(INJURIES_FATAL) as fatalities
FROM crashes WHERE WEATHER_CONDITION IS NOT NULL AND WEATHER_CONDITION != 'UNKNOWN'
GROUP BY WEATHER_CONDITION ORDER BY total_crashes DESC
"""
result3 = pd.read_sql_query(query3, conn)
print("\nQuery 3 — Weather Condition Impact:")
print(result3.to_string(index=False))

# Query 4: Top 15 Contributory Causes
query4 = """
SELECT PRIM_CONTRIBUTORY_CAUSE, COUNT(*) as total_crashes,
       SUM(HAS_INJURY) as injury_crashes,
       ROUND(AVG(HAS_INJURY)*100, 1) as injury_rate_pct,
       SUM(INJURIES_FATAL) as fatalities
FROM crashes WHERE PRIM_CONTRIBUTORY_CAUSE NOT IN ('UNABLE TO DETERMINE', 'NOT APPLICABLE')
GROUP BY PRIM_CONTRIBUTORY_CAUSE ORDER BY total_crashes DESC LIMIT 15
"""
result4 = pd.read_sql_query(query4, conn)
print("\nQuery 4 — Top 15 Contributory Causes:")
print(result4.to_string(index=False))

# Query 5: Hit-and-Run Overview
query5 = """
SELECT HIT_AND_RUN_I, COUNT(*) as total_crashes,
       SUM(HAS_INJURY) as injury_crashes,
       ROUND(AVG(HAS_INJURY)*100, 1) as injury_rate_pct,
       ROUND(AVG(INJURIES_TOTAL), 3) as avg_injuries_per_crash,
       SUM(INJURIES_FATAL) as fatalities
FROM crashes GROUP BY HIT_AND_RUN_I
"""
result5 = pd.read_sql_query(query5, conn)
print("\nQuery 5 — Hit-and-Run Overview:")
print(result5.to_string(index=False))

# Query 6: Top 20 Streets
query6 = """
SELECT STREET_NAME, COUNT(*) as total_crashes,
       SUM(HAS_INJURY) as injury_crashes,
       SUM(INJURIES_FATAL) as fatalities,
       ROUND(AVG(HAS_INJURY)*100, 1) as injury_rate_pct
FROM crashes WHERE STREET_NAME IS NOT NULL
GROUP BY STREET_NAME ORDER BY total_crashes DESC LIMIT 20
"""
result6 = pd.read_sql_query(query6, conn)
print("\nQuery 6 — Top 20 Streets by Crash Volume:")
print(result6.to_string(index=False))

# Key Observations
print("\n" + "=" * 70)
print("KEY OBSERVATIONS FROM SQL EXPLORATION")
print("=" * 70)
total = len(df)
print(f"1. DATASET SCOPE: {total:,} crashes from 2016-2023 in Chicago")
print(f"   - Overall injury rate: {df['HAS_INJURY'].mean()*100:.1f}%")
print(f"   - Overall fatality rate: {df['HAS_FATAL'].mean()*100:.3f}%")
hr_y = df[df['HIT_AND_RUN_I'] == 'Y']
print(f"2. HIT-AND-RUN: {len(hr_y):,} crashes ({len(hr_y)/total*100:.1f}%)")
print(f"3. WET WEATHER: {df['IS_WET_WEATHER'].mean()*100:.1f}% of crashes in wet conditions")

# ============================================================
# PHASE 3: Statistical Analysis
# ============================================================
print("\n" + "=" * 70)
print("PHASE 3: STATISTICAL ANALYSIS")
print("=" * 70)

# --- 3a: Temporal Correlation ---
print("\n--- 3a: TEMPORAL CORRELATION ANALYSIS ---")

# Monthly distribution (expected frequencies proportional to days per month)
monthly_counts = df.groupby('CRASH_MONTH').size()
days_per_month = np.array([31, 28.25, 31, 30, 31, 30, 31, 31, 30, 31, 30, 31])  # avg for leap years
expected_monthly = len(df) * days_per_month / days_per_month.sum()
chi2_month, p_month = stats.chisquare(monthly_counts.values, expected_monthly)
print(f"\n1. Monthly Chi-squared: chi2={chi2_month:,.1f}, p={p_month:.2e}")

# Day-of-week
dow_counts = df.groupby('CRASH_DAY_OF_WEEK').size()
expected_dow = np.full(7, len(df) / 7)
chi2_dow, p_dow = stats.chisquare(dow_counts.values, expected_dow)
print(f"2. Day-of-Week Chi-squared: chi2={chi2_dow:,.1f}, p={p_dow:.2e}")

# Hourly
hourly_counts = df.groupby('CRASH_HOUR').size()
expected_hourly = np.full(24, len(df) / 24)
chi2_hour, p_hour = stats.chisquare(hourly_counts.values, expected_hourly)
print(f"3. Hourly Chi-squared: chi2={chi2_hour:,.1f}, p={p_hour:.2e}")

# Weekend vs weekday (Mann-Whitney U)
df['CRASH_DATE_ONLY'] = df['CRASH_DATE'].dt.date
daily_counts = df.groupby(['CRASH_DATE_ONLY', 'IS_WEEKEND']).size().reset_index(name='count')
weekday_daily = daily_counts[daily_counts['IS_WEEKEND'] == 0]['count']
weekend_daily = daily_counts[daily_counts['IS_WEEKEND'] == 1]['count']
u_stat, p_mw = stats.mannwhitneyu(weekday_daily, weekend_daily, alternative='two-sided')
print(f"4. Weekend vs Weekday Mann-Whitney: U={u_stat:,.0f}, p={p_mw:.2e}")
print(f"   Weekday mean: {weekday_daily.mean():.1f}, Weekend mean: {weekend_daily.mean():.1f}")

# Year-over-year trend
yearly_counts = df.groupby('CRASH_YEAR').size()
r_year, p_year = stats.pearsonr(yearly_counts.index, yearly_counts.values)
print(f"5. Year-over-Year Pearson: r={r_year:.3f}, p={p_year:.4f}")

# Peaks
peak_month = monthly_counts.idxmax()
peak_day = dow_counts.idxmax()
peak_hour = hourly_counts.idxmax()
print(f"\nPeak month: {calendar.month_name[peak_month]} ({monthly_counts.max():,})")
print(f"Peak day: {day_map[peak_day]} ({dow_counts.max():,})")
print(f"Peak hour: {int(peak_hour):02d}:00 ({hourly_counts.max():,})")

# Time series decomposition
print("\nRunning seasonal decomposition...")
monthly_ts = df.groupby([df['CRASH_DATE'].dt.to_period('M')]).size()
monthly_ts.index = monthly_ts.index.to_timestamp()
monthly_ts = monthly_ts.sort_index()
full_range = pd.date_range(monthly_ts.index.min(), monthly_ts.index.max(), freq='MS')
monthly_ts = monthly_ts.reindex(full_range, fill_value=0)
decomp = seasonal_decompose(monthly_ts, model='additive', period=12)

# Chart 1: Seasonal Decomposition — 4-panel (observed/trend/seasonal/residual)
import matplotlib.dates as mdates

fig, axes = plt.subplots(4, 1, figsize=(14, 12), sharex=True)
fig.patch.set_linewidth(0)

covid_date = pd.Timestamp('2020-03-15')

# Panel 1: Observed
axes[0].plot(decomp.observed.index, decomp.observed.values, color='#005B96', linewidth=1.2)
axes[0].set_ylabel('Observed', fontsize=11)
axes[0].set_title('Seasonal Decomposition of Monthly Crash Counts (2016\u20132023)', fontsize=15, fontweight='bold', pad=12)
axes[0].yaxis.set_major_formatter(ticker.FuncFormatter(lambda x, _: f'{x:,.0f}'))
axes[0].yaxis.grid(True, alpha=0.3, linestyle='--')
axes[0].axvspan(pd.Timestamp('2020-03-01'), pd.Timestamp('2020-06-01'), alpha=0.15, color='red', zorder=0)
axes[0].annotate('COVID-19', xy=(covid_date, decomp.observed.min()),
                xytext=(pd.Timestamp('2020-08-01'), decomp.observed.max() * 0.8),
                fontsize=9, color='red', fontweight='bold',
                arrowprops=dict(arrowstyle='->', color='red', lw=1.5))

# Panel 2: Trend
axes[1].plot(decomp.trend.index, decomp.trend.values, color='#E63946', linewidth=2)
axes[1].set_ylabel('Trend', fontsize=11)
axes[1].yaxis.set_major_formatter(ticker.FuncFormatter(lambda x, _: f'{x:,.0f}'))
axes[1].yaxis.grid(True, alpha=0.3, linestyle='--')

# Panel 3: Seasonal
axes[2].plot(decomp.seasonal.index, decomp.seasonal.values, color='#2A9D8F', linewidth=1.2)
axes[2].set_ylabel('Seasonal', fontsize=11)
axes[2].axhline(y=0, color='gray', linestyle='-', linewidth=0.5, alpha=0.5)
axes[2].yaxis.set_major_formatter(ticker.FuncFormatter(lambda x, _: f'{x:,.0f}'))
axes[2].yaxis.grid(True, alpha=0.3, linestyle='--')

# Panel 4: Residual
axes[3].plot(decomp.resid.index, decomp.resid.values, color='#457B9D', linewidth=0.8, alpha=0.7)
axes[3].scatter(decomp.resid.index, decomp.resid.values, s=8, color='#457B9D', alpha=0.5, zorder=3)
axes[3].set_ylabel('Residual', fontsize=11)
axes[3].axhline(y=0, color='gray', linestyle='-', linewidth=0.5, alpha=0.5)
axes[3].set_xlabel('Date', fontsize=11)
axes[3].yaxis.set_major_formatter(ticker.FuncFormatter(lambda x, _: f'{x:,.0f}'))
axes[3].yaxis.grid(True, alpha=0.3, linestyle='--')

for ax_i in axes:
    ax_i.xaxis.grid(False)

fig.text(0.5, -0.01,
         'Additive decomposition (period=12 months). Trend shows COVID-19 dip and recovery. '
         'Seasonal cycle peaks in October, troughs in February. '
         'Large 2020 residual reflects the sudden COVID-19 shock, which the 12-month trend smoother cannot fully absorb.',
         ha='center', fontsize=10, color='gray', style='italic')
plt.tight_layout(rect=[0, 0.03, 1, 1])
plt.subplots_adjust(hspace=0.15)
plt.savefig(f'{OUTPUT_DIR}/01_time_series_decomposition.png', bbox_inches='tight')
plt.close()
print("Saved: 01_time_series_decomposition.png")

# Also compute monthly average for use in slides (Fix 11: exclude partial 2016)
df_full_years = df[df['CRASH_YEAR'] >= 2017]
avg_by_month = df_full_years.groupby(df_full_years['CRASH_DATE'].dt.month).size() / df_full_years['CRASH_YEAR'].nunique()
trough_month = avg_by_month.idxmin()

# --- 3b: Wet Weather Analysis ---
print("\n--- 3b: WET WEATHER & CRASH ANALYSIS ---")

# Filter to known weather only for all weather analysis
df_weather = df[df['WEATHER_KNOWN']].copy()
print(f"Weather analysis: using {len(df_weather):,} crashes with known weather "
      f"(excluded {(~df['WEATHER_KNOWN']).sum():,} NULL/UNKNOWN)")

# Level 1: Naive comparison
wet_inj_rate = df_weather[df_weather['IS_WET_WEATHER'] == 1]['HAS_INJURY'].mean()
dry_inj_rate = df_weather[df_weather['IS_WET_WEATHER'] == 0]['HAS_INJURY'].mean()
print(f"\nLevel 1 — Naive Comparison:")
print(f"  Injury rate in wet weather: {wet_inj_rate*100:.1f}%")
print(f"  Injury rate in dry weather: {dry_inj_rate*100:.1f}%")
print(f"  Difference: {(wet_inj_rate - dry_inj_rate)*100:.1f} pp")

# Level 2: Exposure-adjusted using daily crash counts
# Count crashes per day, grouped by wet/dry
df_weather['CRASH_DATE_ONLY_W'] = df_weather['CRASH_DATE'].dt.date
daily_weather = df_weather.groupby(['CRASH_DATE_ONLY_W', 'IS_WET_WEATHER']).size().reset_index(name='daily_count')
wet_daily = daily_weather[daily_weather['IS_WET_WEATHER'] == 1]['daily_count']
dry_daily = daily_weather[daily_weather['IS_WET_WEATHER'] == 0]['daily_count']
avg_crashes_wet_day = wet_daily.mean()
avg_crashes_dry_day = dry_daily.mean()
cri_wet = avg_crashes_wet_day / avg_crashes_dry_day
n_wet_days = len(wet_daily)
n_dry_days = len(dry_daily)
pct_wet_crashes = df_weather['IS_WET_WEATHER'].mean() * 100
pct_dry_crashes = 100 - pct_wet_crashes
print(f"\nLevel 2 — Exposure-Adjusted (daily crash counts):")
print(f"  Wet days: {n_wet_days:,} days, avg {avg_crashes_wet_day:.1f} crashes/day")
print(f"  Dry days: {n_dry_days:,} days, avg {avg_crashes_dry_day:.1f} crashes/day")
print(f"  CRI (wet vs dry) = {cri_wet:.2f} (ratio of avg daily crash counts)")
print(f"  NOTE: 'Wet days' are classified by whether any crash on that day had a wet-weather code,")
print(f"  not by independent weather station data. Many crashes on a 'wet day' likely occurred during")
print(f"  dry hours. This CRI should be interpreted with caution.")

# Level 3: Chi-squared test (using known-weather subset)
contingency = pd.crosstab(df_weather['IS_WET_WEATHER'], df_weather['HAS_INJURY'])
chi2_weather, p_chi_weather, dof_w, expected_w = stats.chi2_contingency(contingency)
assert (expected_w >= 5).all(), f"Expected frequency check failed: min={expected_w.min():.1f}"
print(f"\nLevel 3 — Chi-squared test of independence:")
print(f"  chi2={chi2_weather:,.1f}, df={dof_w}, p={p_chi_weather:.2e}")

# Level 4: Odds ratio
a = contingency.iloc[1, 1]  # Wet & Injury
b = contingency.iloc[1, 0]  # Wet & No Injury
c = contingency.iloc[0, 1]  # Dry & Injury
d = contingency.iloc[0, 0]  # Dry & No Injury
OR = (a * d) / (b * c)
log_OR = np.log(OR)
SE_log_OR = np.sqrt(1/a + 1/b + 1/c + 1/d)
CI_lower = np.exp(log_OR - 1.96 * SE_log_OR)
CI_upper = np.exp(log_OR + 1.96 * SE_log_OR)
print(f"\nLevel 4 — Unadjusted Odds Ratio:")
print(f"  OR={OR:.3f}, 95% CI=({CI_lower:.3f}, {CI_upper:.3f})")

# Level 5: Logistic regression
print("\nLevel 5 — Logistic Regression (adjusted for confounders incl. month)...")
lr_df = df_weather[['HAS_INJURY', 'IS_WET_WEATHER', 'IS_DARK', 'POSTED_SPEED_LIMIT',
            'CRASH_HOUR', 'CRASH_DAY_OF_WEEK', 'CRASH_MONTH']].dropna()
model = smf.logit('HAS_INJURY ~ IS_WET_WEATHER + IS_DARK + POSTED_SPEED_LIMIT + C(CRASH_HOUR) + C(CRASH_DAY_OF_WEEK) + C(CRASH_MONTH)',
                   data=lr_df)
result = model.fit(maxiter=200, method='lbfgs', disp=False)

key_vars = ['IS_WET_WEATHER', 'IS_DARK', 'POSTED_SPEED_LIMIT']
print(f"  Model converged: {result.mle_retvals['converged']}")
print(f"  Pseudo R-squared: {result.prsquared:.4f}")
print(f"  Observations: {int(result.nobs):,}")

or_data = []
for var in key_vars:
    coef = result.params[var]
    ci = result.conf_int().loc[var]
    or_val = np.exp(coef)
    or_lower = np.exp(ci.iloc[0])
    or_upper = np.exp(ci.iloc[1])
    p = result.pvalues[var]
    print(f"  {var}: OR={or_val:.3f}, 95% CI=({or_lower:.3f}, {or_upper:.3f}), p={p:.2e}")
    or_data.append({'Variable': var, 'OR': or_val, 'CI_lower': or_lower, 'CI_upper': or_upper, 'p': p})
or_df = pd.DataFrame(or_data)

wet_or = or_df[or_df['Variable'] == 'IS_WET_WEATHER']['OR'].values[0]
print(f"\n  -> Adjusted wet weather OR = {wet_or:.3f}")

# Level 6: Stratified analysis (using known-weather subset)
print("\nLevel 6 — Stratified Analysis by Lighting:")
strata = df_weather.groupby('LIGHTING_CONDITION')
for name, group in strata:
    ct = pd.crosstab(group['IS_WET_WEATHER'], group['HAS_INJURY'])
    if ct.shape == (2, 2) and (ct > 0).all().all():
        a_s = ct.iloc[1, 1]; b_s = ct.iloc[1, 0]
        c_s = ct.iloc[0, 1]; d_s = ct.iloc[0, 0]
        or_s = (a_s * d_s) / (b_s * c_s)
        se_s = np.sqrt(1/a_s + 1/b_s + 1/c_s + 1/d_s)
        ci_l_s = np.exp(np.log(or_s) - 1.96 * se_s)
        ci_u_s = np.exp(np.log(or_s) + 1.96 * se_s)
        print(f"  {name:<30} N={len(group):<10,} OR={or_s:.3f} ({ci_l_s:.3f}, {ci_u_s:.3f})")

# --- 3c: Key Findings ---
print("\n--- 3c: KEY FINDINGS ---")

# Finding 1: Hit-and-Run
print("\nFINDING 1: HIT-AND-RUN PATTERNS")
hr = df[df['HIT_AND_RUN_I'] == 'Y']
non_hr = df[df['HIT_AND_RUN_I'] == 'N']
print(f"  Total H&R: {len(hr):,} ({len(hr)/len(df)*100:.1f}%)")
print(f"  H&R injury rate: {hr['HAS_INJURY'].mean()*100:.1f}%")
print(f"  Non-H&R injury rate: {non_hr['HAS_INJURY'].mean()*100:.1f}%")
print(f"  H&R fatality rate: {hr['HAS_FATAL'].mean()*100:.3f}%")
print(f"  Non-H&R fatality rate: {non_hr['HAS_FATAL'].mean()*100:.3f}%")

hr_contingency = pd.crosstab(df['HIT_AND_RUN_I'], df['HAS_INJURY'])
chi2_hr, p_hr, _, _ = stats.chi2_contingency(hr_contingency)
print(f"  Chi-squared (H&R x Injury): chi2={chi2_hr:,.1f}, p={p_hr:.2e}")

hr_hourly = hr.groupby('CRASH_HOUR').size() / df.groupby('CRASH_HOUR').size() * 100
print(f"  Peak H&R hours:")
for h, pct in hr_hourly.nlargest(5).items():
    print(f"    {int(h):02d}:00 - {pct:.1f}%")

hr_rate_hour = df.groupby('CRASH_HOUR')['HIT_AND_RUN_I'].apply(lambda x: (x == 'Y').mean() * 100)
hr_rate_dow = df.groupby('DAY_NAME')['HIT_AND_RUN_I'].apply(lambda x: (x == 'Y').mean() * 100)

# Finding 2: Speed & Severity
print("\nFINDING 2: SPEED LIMIT & CRASH SEVERITY")
speed_order = ['0-20 mph', '21-30 mph', '31-40 mph', '45+ mph']
speed_stats = df[df['SPEED_CATEGORY'] != 'Unknown'].groupby('SPEED_CATEGORY').agg(
    total_crashes=('HAS_INJURY', 'count'),
    injury_rate=('HAS_INJURY', 'mean'),
    fatality_rate=('HAS_FATAL', 'mean'),
    avg_injuries=('INJURIES_TOTAL', 'mean')
).reindex(speed_order)
speed_stats['injury_rate_pct'] = speed_stats['injury_rate'] * 100
speed_stats['fatality_rate_per_1000'] = speed_stats['fatality_rate'] * 1000

for idx, row in speed_stats.iterrows():
    print(f"  {idx}: {int(row['total_crashes']):,} crashes, "
          f"injury rate={row['injury_rate_pct']:.1f}%, "
          f"fatal/1000={row['fatality_rate_per_1000']:.2f}")

valid = df[df['POSTED_SPEED_LIMIT'].notna() & df['INJURIES_TOTAL'].notna()]
rho, p_spearman = stats.spearmanr(valid['POSTED_SPEED_LIMIT'], valid['INJURIES_TOTAL'])
rpb, p_pb = stats.pointbiserialr(valid['HAS_FATAL'], valid['POSTED_SPEED_LIMIT'])
print(f"  Spearman (speed vs injuries): rho={rho:.3f}, p={p_spearman:.2e}")
print(f"  Point-biserial (speed vs fatal): r={rpb:.3f}, p={p_pb:.2e}")

# Finding 3: Dangerous Causes
print("\nFINDING 3: MOST DANGEROUS CONTRIBUTORY CAUSES")
causes_df = df[~df['PRIM_CONTRIBUTORY_CAUSE'].isin(['UNABLE TO DETERMINE', 'NOT APPLICABLE'])].copy()
cause_stats = causes_df.groupby('PRIM_CONTRIBUTORY_CAUSE').agg(
    total_crashes=('HAS_INJURY', 'count'),
    injury_rate=('HAS_INJURY', 'mean'),
    fatalities=('INJURIES_FATAL', 'sum')
).reset_index()
cause_stats['injury_crash_count'] = cause_stats['total_crashes'] * cause_stats['injury_rate']
cause_stats = cause_stats.sort_values('injury_crash_count', ascending=False).head(15)
cause_stats['injury_rate_pct'] = cause_stats['injury_rate'] * 100

for i, (_, row) in enumerate(cause_stats.iterrows(), 1):
    print(f"  {i}. {row['PRIM_CONTRIBUTORY_CAUSE'][:45]}: "
          f"{int(row['total_crashes']):,} crashes, "
          f"injury rate={row['injury_rate_pct']:.1f}%")

# --- 3d: Hotspot Identification ---
print("\n--- 3d: CRASH HOTSPOT IDENTIFICATION ---")
injury_crashes = df[(df['HAS_INJURY'] == 1) &
                    df['LATITUDE'].notna() & df['LONGITUDE'].notna() &
                    df['LATITUDE'].between(41.6, 42.1) &
                    df['LONGITUDE'].between(-87.95, -87.5)].copy()
print(f"Injury crashes with valid coordinates: {len(injury_crashes):,}")

injury_sample = injury_crashes  # Use all injury crashes
eps_m = 150  # metres
min_samp = 100
coords_rad = np.radians(injury_sample[['LATITUDE', 'LONGITUDE']].values)
print(f"Running DBSCAN on {len(injury_sample):,} injury crashes (eps={eps_m}m, min_samples={min_samp}, haversine)...")
db = DBSCAN(eps=eps_m/6371000, min_samples=min_samp, metric='haversine', algorithm='ball_tree')
injury_sample = injury_sample.copy()
injury_sample['cluster'] = db.fit_predict(coords_rad)

n_clusters = len(set(injury_sample['cluster'])) - (1 if -1 in injury_sample['cluster'] else 0)
n_noise = (injury_sample['cluster'] == -1).sum()
print(f"Clusters found: {n_clusters}")
print(f"Noise points: {n_noise:,} ({n_noise/len(injury_sample)*100:.1f}%)")

# Profile top 20 clusters
cluster_profiles = []
for cid in injury_sample[injury_sample['cluster'] >= 0]['cluster'].value_counts().head(20).index:
    cluster = injury_sample[injury_sample['cluster'] == cid]
    profile = {
        'cluster_id': cid,
        'crash_count': len(cluster),
        'center_lat': cluster['LATITUDE'].mean(),
        'center_lon': cluster['LONGITUDE'].mean(),
        'total_injuries': cluster['INJURIES_TOTAL'].sum(),
        'fatalities': cluster['INJURIES_FATAL'].sum(),
        'top_cause': cluster['PRIM_CONTRIBUTORY_CAUSE'].mode().iloc[0] if len(cluster) > 0 else 'N/A',
        'top_lighting': cluster['LIGHTING_CONDITION'].mode().iloc[0] if len(cluster) > 0 else 'N/A',
        'mean_speed_limit': cluster['POSTED_SPEED_LIMIT'].mean(),
        'top_street': cluster['STREET_NAME'].mode().iloc[0] if len(cluster['STREET_NAME'].dropna()) > 0 else 'N/A'
    }
    cluster_profiles.append(profile)

hotspots_df = pd.DataFrame(cluster_profiles)
print("\nTop 20 Crash Hotspots:")
for i, row in hotspots_df.iterrows():
    print(f"  {i+1}. {row['top_street']}: {int(row['crash_count'])} crashes, "
          f"{int(row['total_injuries'])} injuries, {int(row['fatalities'])} fatal")

# Compute summary variables early (used in charts and slides)
total_crashes = len(df)
total_injuries = int(df['INJURIES_TOTAL'].sum())
total_fatalities = int(df['INJURIES_FATAL'].sum())
hr = df[df['HIT_AND_RUN_I'] == 'Y']
hr_pct = len(hr) / len(df) * 100

# Dynamic references for slide text
top_hotspot_street = hotspots_df.iloc[0]['top_street'] if len(hotspots_df) > 0 else 'N/A'
top_4_streets = ', '.join(hotspots_df.head(4)['top_street'].tolist()) if len(hotspots_df) >= 4 else top_hotspot_street

# Peak/trough computations for slide text
peak_month_name = calendar.month_name[peak_month]
trough_month_name = calendar.month_name[trough_month]
peak_day_name = day_map[peak_day]
peak_hour_label_main = f"{int(peak_hour) % 12 or 12} {'AM' if peak_hour < 12 else 'PM'}"
injury_rate = df['HAS_INJURY'].mean() * 100
injury_ratio = int(round(1 / (injury_rate / 100)))  # e.g., 7 for ~14%

# Weather injury rates for slide text
clear_inj_rate = df_weather[df_weather['WEATHER_CONDITION'] == 'CLEAR']['HAS_INJURY'].mean() * 100
rain_inj_rate = df_weather[df_weather['WEATHER_CONDITION'] == 'RAIN']['HAS_INJURY'].mean() * 100

# Safest hour for slide text
safest_hour = hourly_counts.idxmin()
safest_label = f"{int(safest_hour) % 12 or 12}–{int((safest_hour+2) % 24) % 12 or 12} {'AM' if safest_hour < 12 else 'PM'}"

# ============================================================
# PHASE 4: VISUALIZATIONS
# ============================================================
print("\n" + "=" * 70)
print("PHASE 4: VISUALIZATION GENERATION")
print("=" * 70)

# Chart 2: Heatmap
fig, ax = plt.subplots(figsize=(14, 6))
fig.patch.set_linewidth(0)
sns.heatmap(heatmap_data, cmap='YlOrRd', annot=False, ax=ax,
            linewidths=0.5, linecolor='white',
            cbar_kws={'label': 'Number of Crashes'})
ax.set_title('Crash Density by Day of Week and Hour', fontsize=14, fontweight='bold')
fig.text(0.5, 0.94, f'Darker = more crashes. Rush hours and {peak_day_name} afternoons are highest risk',
         ha='center', fontsize=10, color='gray', style='italic')
ax.set_xlabel('Hour of Day')
ax.set_ylabel('')
# AM/PM x-tick labels every 3 hours
hour_labels = ['12am', '', '', '3am', '', '', '6am', '', '', '9am', '', '',
               '12pm', '', '', '3pm', '', '', '6pm', '', '', '9pm', '', '']
ax.set_xticks(np.arange(24) + 0.5)
ax.set_xticklabels(hour_labels, fontsize=9)
# Peak annotation
peak_val = heatmap_data.max().max()
for day in heatmap_data.index:
    for hour in heatmap_data.columns:
        if heatmap_data.loc[day, hour] == peak_val:
            day_idx = list(heatmap_data.index).index(day)
            hour_idx = list(heatmap_data.columns).index(hour)
            rect = plt.Rectangle((hour_idx, day_idx), 1, 1, fill=False, edgecolor='red', linewidth=3)
            ax.add_patch(rect)
            ax.annotate(f'Peak: {day} {int(hour)}:00',
                       xy=(hour_idx + 0.5, day_idx + 0.5),
                       xytext=(hour_idx + 3, max(0, day_idx - 1.5)),
                       fontsize=10, fontweight='bold', color='red',
                       arrowprops=dict(arrowstyle='->', color='red', lw=1.5))
            break
    else:
        continue
    break
plt.tight_layout(rect=[0, 0, 1, 0.93])
plt.savefig(f'{OUTPUT_DIR}/02_dow_hour_heatmap.png', bbox_inches='tight')
plt.close()
print("Saved: 02_dow_hour_heatmap.png")

# Chart 3: Monthly bar
monthly_by_year = df_full_years.groupby(['CRASH_YEAR', 'CRASH_MONTH']).size().reset_index(name='count')
monthly_avg = monthly_by_year.groupby('CRASH_MONTH')['count'].agg(['mean', 'std']).reset_index()
month_names = [calendar.month_abbr[int(m)] for m in monthly_avg['CRASH_MONTH']]

fig, ax = plt.subplots(figsize=(12, 5))
fig.patch.set_linewidth(0)
bars = ax.bar(month_names, monthly_avg['mean'], yerr=monthly_avg['std'],
              capsize=4, color='#005B96', alpha=0.85, edgecolor='white')
ax.set_title('Average Monthly Crash Count (2017–2023) with Year-to-Year Variation',
             fontsize=14, fontweight='bold')
fig.text(0.5, 0.93, f'{peak_month_name} sees highest crash volume; summer months also elevated',
         ha='center', fontsize=10, color='gray', style='italic')
ax.set_ylabel('Crash Count')
ax.set_xlabel('Month')
ax.yaxis.grid(True, alpha=0.3, linestyle='--')
ax.yaxis.set_major_formatter(ticker.FuncFormatter(lambda x, _: f'{x:,.0f}'))
peak_idx = monthly_avg['mean'].idxmax()
bars[peak_idx].set_color('#E63946')
# Data labels on top of each bar
for i, (val, std_val) in enumerate(zip(monthly_avg['mean'], monthly_avg['std'])):
    ax.text(i, val + std_val + 100, f'{val:,.0f}', ha='center', va='bottom', fontsize=8, fontweight='bold')
# Peak annotation
ax.annotate(f'Peak: {peak_month_name}', xy=(peak_idx, monthly_avg['mean'].iloc[peak_idx]),
           xytext=(peak_idx - 4, monthly_avg['mean'].iloc[peak_idx] + monthly_avg['std'].iloc[peak_idx] + 800),
           fontsize=10, fontweight='bold', color='#E63946',
           arrowprops=dict(arrowstyle='->', color='#E63946', lw=1.5))
# Error bars explanation
ax.text(0.98, 0.02, 'Error bars = \u00b11 standard deviation across years (year-to-year variation)',
        transform=ax.transAxes, ha='right', va='bottom', fontsize=9, color='#333', style='italic',
        bbox=dict(boxstyle='round,pad=0.3', facecolor='lightyellow', alpha=0.8, edgecolor='none'))
plt.tight_layout(rect=[0, 0, 1, 0.92])
plt.savefig(f'{OUTPUT_DIR}/03_monthly_bar.png', bbox_inches='tight')
plt.close()
print("Saved: 03_monthly_bar.png")

# Chart 4: Weather Impact (single-axis horizontal bar with annotations)
weather_plot = df_weather[~df_weather['WEATHER_CONDITION'].isin(['OTHER'])].groupby('WEATHER_CONDITION').agg(
    total_crashes=('HAS_INJURY', 'count'),
    injury_rate=('HAS_INJURY', 'mean')
).sort_values('total_crashes', ascending=True)
weather_plot['injury_rate_pct'] = weather_plot['injury_rate'] * 100

fig, ax = plt.subplots(figsize=(12, 6))
fig.patch.set_linewidth(0)
y_pos = range(len(weather_plot))
ax.barh(y_pos, weather_plot['total_crashes'], color='#005B96', alpha=0.8, edgecolor='white')
ax.set_yticks(y_pos)
ax.set_yticklabels(weather_plot.index, fontsize=10)
ax.set_xlabel('Total Crashes')
ax.set_title('Crash Count and Injury Rate by Weather Condition',
              fontsize=14, fontweight='bold')
fig.text(0.5, 0.94, 'Rain and freezing conditions show elevated injury rates despite lower crash volumes',
         ha='center', fontsize=10, color='gray', style='italic')
ax.xaxis.set_major_formatter(ticker.FuncFormatter(lambda x, _: f'{x:,.0f}'))
ax.yaxis.grid(False)
ax.xaxis.grid(True, alpha=0.3, linestyle='--')
# Annotate: crash count inside/beside bars, injury rate well to the right
max_crashes_w = weather_plot['total_crashes'].max()
for i, (crashes, inj_rate) in enumerate(zip(weather_plot['total_crashes'], weather_plot['injury_rate_pct'])):
    # Crash count: inside bar if big enough, otherwise just outside
    if crashes > max_crashes_w * 0.15:
        ax.text(crashes * 0.5, i, f'{crashes:,.0f}', ha='center', va='center',
                fontsize=9, fontweight='bold', color='white')
        # Injury rate: positioned at a fixed offset past the bar end
        rate_x = crashes + max_crashes_w * 0.02
    else:
        ax.text(crashes + max_crashes_w * 0.01, i, f'{crashes:,.0f}', ha='left', va='center',
                fontsize=8, color='#333333')
        # For small bars, offset injury rate past the count text
        rate_x = crashes + max_crashes_w * 0.08
    ax.text(rate_x, i + 0.25, f'{inj_rate:.1f}% injury rate',
            ha='left', va='center', fontsize=8, color='#E63946', fontweight='bold')
# Extend x-axis to make room for annotations
ax.set_xlim(0, max_crashes_w * 1.25)
# Footnote about exposure
fig.text(0.5, -0.02,
         'Note: CLEAR weather dominates because it is the most common condition. See exposure-adjusted analysis for context.',
         ha='center', fontsize=8, color='gray', style='italic')
plt.tight_layout(rect=[0, 0.03, 1, 0.93])
plt.savefig(f'{OUTPUT_DIR}/04_weather_impact.png', bbox_inches='tight')
plt.close()
print("Saved: 04_weather_impact.png")

# Chart 5: "How much does each factor increase injury risk?" — intuitive bar chart
nobs = int(result.nobs)
r2 = result.prsquared

# Compute % increase from OR: (OR - 1) * 100
or_df_plot = or_df.copy()
or_df_plot['pct_increase'] = (or_df_plot['OR'] - 1) * 100
or_df_plot['ci_lo_pct'] = (or_df_plot['CI_lower'] - 1) * 100
or_df_plot['ci_hi_pct'] = (or_df_plot['CI_upper'] - 1) * 100

# Human-readable labels
readable_labels = {
    'IS_WET_WEATHER': 'Wet Weather\n(rain, snow, ice)',
    'IS_DARK': 'Darkness\n(no daylight)',
    'POSTED_SPEED_LIMIT': 'Speed Limit\n(each +1 mph)'
}
or_df_plot['label'] = or_df_plot['Variable'].map(readable_labels)

fig, ax = plt.subplots(figsize=(10, 5))
fig.patch.set_linewidth(0)
y_pos5 = list(range(len(or_df_plot)))
bar_colors5 = ['#E63946' if pct > 10 else '#F4A261' if pct > 5 else '#457B9D'
               for pct in or_df_plot['pct_increase']]

ax.barh(y_pos5, or_df_plot['pct_increase'].values, color=bar_colors5, alpha=0.85,
        edgecolor='white', height=0.55)

# Confidence interval lines
for i, (_, row) in enumerate(or_df_plot.iterrows()):
    ax.plot([row['ci_lo_pct'], row['ci_hi_pct']], [i, i],
            color='#333', linewidth=2, solid_capstyle='round')
    ax.plot([row['ci_lo_pct'], row['ci_lo_pct']], [i-0.1, i+0.1], color='#333', linewidth=1.5)
    ax.plot([row['ci_hi_pct'], row['ci_hi_pct']], [i-0.1, i+0.1], color='#333', linewidth=1.5)

# Data labels — plain English
for i, (_, row) in enumerate(or_df_plot.iterrows()):
    pct = row['pct_increase']
    ci_lo = row['ci_lo_pct']
    ci_hi = row['ci_hi_pct']
    ax.text(row['ci_hi_pct'] + 0.5, i,
            f'+{pct:.1f}%  (range: {ci_lo:.1f}\u2013{ci_hi:.1f}%)',
            ha='left', va='center', fontsize=10, fontweight='bold', color='#333')

ax.set_yticks(y_pos5)
ax.set_yticklabels(or_df_plot['label'].values, fontsize=11)
ax.set_xlabel('% Increase in Injury Risk (controlling for other factors)', fontsize=11)
ax.axvline(x=0, color='gray', linestyle='-', linewidth=0.8, alpha=0.5)
ax.yaxis.grid(False)
ax.xaxis.grid(True, alpha=0.3, linestyle='--')
ax.xaxis.set_major_formatter(ticker.FuncFormatter(lambda x, _: f'+{x:.0f}%' if x > 0 else f'{x:.0f}%'))

ax.set_title('How Much Does Each Factor Increase Injury Risk?',
             fontsize=14, fontweight='bold')
fig.text(0.5, 0.93,
         f'Logistic regression controlling for all factors simultaneously (incl. month)  |  '
         f'n = {nobs:,}  |  All p < 0.001  |  Pseudo R\u00b2 = {r2:.3f}  |  Horizontal lines = 95% CI',
         ha='center', fontsize=9, color='gray', style='italic')

# Compounding speed footnote
spd_pct = or_df_plot[or_df_plot['Variable'] == 'POSTED_SPEED_LIMIT']['pct_increase'].values[0]
cumul_25 = ((1 + spd_pct/100)**25 - 1) * 100
fig.text(0.5, -0.02,
         f'Note: Speed effect is multiplicative. Over a 25 mph range (20\u219245 mph), '
         f'the cumulative increase is +{cumul_25:.0f}% injury odds.',
         ha='center', fontsize=8, color='#666', style='italic')

# Extend x-axis for labels
xmax5 = or_df_plot['ci_hi_pct'].max()
ax.set_xlim(-1, xmax5 + 14)

plt.tight_layout(rect=[0, 0.04, 1, 0.92])
plt.savefig(f'{OUTPUT_DIR}/05_odds_ratio_chart.png', bbox_inches='tight')
plt.close()
print("Saved: 05_odds_ratio_chart.png")

# Chart 6: Hit-and-Run Temporal
fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 5))
fig.patch.set_linewidth(0)
overall_hr_rate = (df['HIT_AND_RUN_I'] == 'Y').mean() * 100

# Left: hourly
ax1.plot(hr_rate_hour.index, hr_rate_hour.values, 'o-', color='#E63946', linewidth=2, markersize=5)
ax1.fill_between(hr_rate_hour.index, hr_rate_hour.values, alpha=0.2, color='#E63946')
ax1.axhline(y=overall_hr_rate, color='gray', linestyle='--', linewidth=1, alpha=0.7)
ax1.text(23, overall_hr_rate + 0.5, f'Overall avg: {overall_hr_rate:.1f}%', ha='right', fontsize=8, color='gray')
ax1.set_xlabel('Hour of Day')
ax1.set_ylabel('Hit-and-Run Rate (%)')
ax1.set_title('(a) Hit-and-Run Rate by Hour', fontweight='bold')
hour_ticks = list(range(0, 24, 3))
hour_tick_labels = ['12am', '3am', '6am', '9am', '12pm', '3pm', '6pm', '9pm']
ax1.set_xticks(hour_ticks)
ax1.set_xticklabels(hour_tick_labels, fontsize=9)
ax1.yaxis.grid(True, alpha=0.3, linestyle='--')
# Peak annotation
peak_hr_hour = hr_rate_hour.idxmax()
peak_hr_val = hr_rate_hour.max()
ax1.annotate(f'Peak: {int(peak_hr_hour) % 12 or 12} {"AM" if peak_hr_hour < 12 else "PM"} ({peak_hr_val:.0f}%)',
            xy=(peak_hr_hour, peak_hr_val),
            xytext=(peak_hr_hour + 4, peak_hr_val - 5),
            fontsize=9, fontweight='bold', color='darkred',
            arrowprops=dict(arrowstyle='->', color='darkred', lw=1.5))

# Right: day of week bars
hr_rate_dow_ordered = hr_rate_dow.reindex(day_order)
ax2.bar(range(7), hr_rate_dow_ordered.values, color='#E63946', alpha=0.8, edgecolor='white')
ax2.set_xticks(range(7))
ax2.set_xticklabels([d[:3] for d in day_order], fontsize=10)
ax2.set_ylabel('Hit-and-Run Rate (%)')
ax2.set_title('(b) Hit-and-Run Rate by Day', fontweight='bold')
ax2.yaxis.grid(True, alpha=0.3, linestyle='--')
ax2.axhline(y=overall_hr_rate, color='gray', linestyle='--', linewidth=1, alpha=0.7)
# Data labels on bars
for i, val in enumerate(hr_rate_dow_ordered.values):
    ax2.text(i, val + 0.2, f'{val:.1f}%', ha='center', va='bottom', fontsize=8, fontweight='bold')

fig.suptitle('Hit-and-Run Crash Patterns', fontsize=14, fontweight='bold', y=1.02)
fig.text(0.5, 0.98, 'Late-night hours (midnight\u20134 AM) have elevated hit-and-run rates \u2014 possibly linked to impairment or reduced witnesses',
         ha='center', fontsize=10, color='gray', style='italic')
plt.tight_layout()
plt.savefig(f'{OUTPUT_DIR}/06_hit_and_run.png', bbox_inches='tight')
plt.close()
print("Saved: 06_hit_and_run.png")

# Chart 7: Speed vs Severity (dual y-axis: bars for injury rate, line for fatality rate)
fig, ax1 = plt.subplots(figsize=(10, 5))
fig.patch.set_linewidth(0)
x = np.arange(len(speed_order))
# Bars for injury rate on left axis
bars7 = ax1.bar(x, speed_stats['injury_rate_pct'], width=0.5, color='#E63946', alpha=0.85,
               edgecolor='white', label='Injury Rate (%)', zorder=3)
ax1.set_ylabel('Injury Rate (%)', color='#E63946', fontsize=12)
ax1.tick_params(axis='y', labelcolor='#E63946')
ax1.yaxis.grid(True, alpha=0.3, linestyle='--')
# Line for fatality rate on right axis
ax2 = ax1.twinx()
ax2.plot(x, speed_stats['fatality_rate_per_1000'], 's-', color='#1D3557',
         linewidth=2.5, markersize=10, markerfacecolor='white', markeredgewidth=2,
         label='Fatality Rate (per 1,000)', zorder=4)
ax2.set_ylabel('Fatality Rate (per 1,000 crashes)', color='#1D3557', fontsize=12)
ax2.tick_params(axis='y', labelcolor='#1D3557')
# X-axis labels with sample sizes
x_labels = [f'{cat}\n(n={int(speed_stats.loc[cat, "total_crashes"]):,})' for cat in speed_order]
ax1.set_xticks(x)
ax1.set_xticklabels(x_labels, fontsize=10)
ax1.set_xlabel('Posted Speed Limit')
ax1.set_title('Crash Severity by Posted Speed Limit', fontsize=14, fontweight='bold')
fig.text(0.5, 0.93, 'Higher speed limits are strongly associated with more severe crash outcomes',
         ha='center', fontsize=10, color='gray', style='italic')
# Data labels
for i, (inj, fat) in enumerate(zip(speed_stats['injury_rate_pct'], speed_stats['fatality_rate_per_1000'])):
    ax1.text(i, inj + 0.5, f'{inj:.1f}%', ha='center', fontsize=9, fontweight='bold', color='#E63946')
    ax2.text(i, fat + 0.15, f'{fat:.1f}', ha='center', fontsize=9, fontweight='bold', color='#1D3557')
# Dose-response callout
fat_min = speed_stats['fatality_rate_per_1000'].iloc[0]
fat_max = speed_stats['fatality_rate_per_1000'].iloc[-1]
multiplier = fat_max / fat_min if fat_min > 0 else 0
ax1.text(0.98, 0.02, f'{multiplier:.1f}x increase in fatality rate\nfrom lowest to highest speed zone',
         transform=ax1.transAxes, ha='right', va='bottom', fontsize=9,
         fontweight='bold', color='#1D3557',
         bbox=dict(boxstyle='round,pad=0.3', facecolor='lightyellow', alpha=0.8))
# Combined legend
lines1, labels1 = ax1.get_legend_handles_labels()
lines2, labels2 = ax2.get_legend_handles_labels()
ax1.legend(lines1 + lines2, labels1 + labels2, loc='upper left')
plt.tight_layout(rect=[0, 0, 1, 0.92])
plt.savefig(f'{OUTPUT_DIR}/07_speed_severity.png', bbox_inches='tight')
plt.close()
print("Saved: 07_speed_severity.png")

# Chart 8: Top Contributory Causes
fig, ax = plt.subplots(figsize=(13, 7))
fig.patch.set_linewidth(0)
cs = cause_stats.sort_values('injury_crash_count', ascending=True)
norm = plt.Normalize(cs['injury_rate_pct'].min(), cs['injury_rate_pct'].max())
colors_cs = plt.cm.RdYlGn_r(norm(cs['injury_rate_pct'].values))
bars8 = ax.barh(range(len(cs)), cs['total_crashes'], color=colors_cs, edgecolor='white', height=0.7)
ax.set_yticks(range(len(cs)))
ax.set_yticklabels(cs['PRIM_CONTRIBUTORY_CAUSE'].str.title().str[:60], fontsize=9)
ax.set_xlabel('Total Crashes')
ax.set_title('Top 15 Contributory Causes (Ranked by Injury Volume, Bars = Total Crashes)', fontsize=14, fontweight='bold')
fig.text(0.5, 0.95, 'Ranked by total injury crashes (= count × injury rate). Bar length = total crashes. Color = injury rate.',
         ha='center', fontsize=10, color='gray', style='italic')
ax.xaxis.set_major_formatter(ticker.FuncFormatter(lambda x, _: f'{x:,.0f}'))
ax.yaxis.grid(False)
ax.xaxis.grid(True, alpha=0.3, linestyle='--')
# Crash count inside bars & injury rate to the right
max_crashes_c = cs['total_crashes'].max()
for i, (crashes, inj_rate) in enumerate(zip(cs['total_crashes'], cs['injury_rate_pct'])):
    text_color = 'white' if crashes > max_crashes_c * 0.3 else 'black'
    ax.text(min(crashes * 0.5, crashes - max_crashes_c * 0.02), i, f'{crashes:,.0f}',
            ha='center', va='center', fontsize=8, fontweight='bold', color=text_color)
    ax.text(crashes + max_crashes_c * 0.01, i, f'{inj_rate:.1f}% injury',
            va='center', fontsize=8, color='gray')
# Top 3 bold border
for j in range(max(0, len(cs) - 3), len(cs)):
    bars8[j].set_edgecolor('black')
    bars8[j].set_linewidth(2)
sm_cmap = plt.cm.ScalarMappable(cmap='RdYlGn_r', norm=norm)
sm_cmap.set_array([])
cbar = plt.colorbar(sm_cmap, ax=ax, shrink=0.6)
cbar.set_label('Injury Rate (%) \u2014 Green: Lower, Red: Higher', fontsize=9)
plt.tight_layout(rect=[0, 0, 1, 0.94])
plt.savefig(f'{OUTPUT_DIR}/08_contributory_causes.png', bbox_inches='tight')
plt.close()
print("Saved: 08_contributory_causes.png")

# Chart 9: Static Crash Density Map
fig, ax = plt.subplots(figsize=(10, 12))
fig.patch.set_linewidth(0)
valid_coords = df[df['LATITUDE'].notna() & df['LONGITUDE'].notna() &
                  df['LATITUDE'].between(41.64, 42.02) &
                  df['LONGITUDE'].between(-87.84, -87.52)]
sample = valid_coords.sample(n=min(50000, len(valid_coords)), random_state=42)
ax.scatter(sample['LONGITUDE'], sample['LATITUDE'], alpha=0.02, s=1, c='#005B96')
ax.set_title('Chicago Crash Density with Identified Hotspots', fontsize=14, fontweight='bold')
fig.text(0.5, 0.97, f'DBSCAN clustering ({eps_m}m radius, min {min_samp} crashes). Red circles = injury hotspots',
         ha='center', fontsize=10, color='gray', style='italic')
if len(hotspots_df) > 0:
    sizes = np.log1p(hotspots_df['crash_count']) / np.log1p(hotspots_df['crash_count'].max()) * 200
    sizes = sizes.clip(upper=200)  # Cap max circle size to avoid obscuring density
    ax.scatter(hotspots_df['center_lon'], hotspots_df['center_lat'],
               s=sizes, c='red', alpha=0.7, edgecolors='darkred', linewidth=1.5,
               zorder=5, label='Crash Hotspots')
    # Numbered hotspot labels with staggered offsets to reduce overlap
    offsets = [(12, 12), (-18, -14), (12, -14), (-18, 12), (16, 8),
               (-20, -10), (12, -18), (-16, 14), (18, -8), (-14, 16)]
    for idx, (i, row) in enumerate(hotspots_df.head(10).iterrows()):
        offset_x, offset_y = offsets[idx % len(offsets)]
        ax.annotate(f"{idx+1}. {row['top_street']}",
                   (row['center_lon'], row['center_lat']),
                   fontsize=9, fontweight='bold', color='darkred',
                   xytext=(offset_x, offset_y), textcoords='offset points',
                   bbox=dict(boxstyle='round,pad=0.2', facecolor='white', alpha=0.7, edgecolor='none'))
# Geographic reference labels
geo_refs = [(-87.63, 41.882, 'Downtown'), (-87.63, 41.76, 'South Side'),
            (-87.75, 41.98, "O'Hare Area"), (-87.72, 41.88, 'West Side')]
for lon, lat, name in geo_refs:
    ax.text(lon, lat, name, fontsize=10, color='navy', alpha=0.4, fontweight='bold', ha='center', style='italic')
# Size legend
for s_val, s_label in [(50, 'Small cluster'), (150, 'Medium cluster'), (300, 'Large cluster')]:
    ax.scatter([], [], s=s_val, c='red', alpha=0.7, edgecolors='darkred', label=s_label)
# Density annotation
ax.text(0.02, 0.02, f'Blue density = all crashes ({min(50000, len(valid_coords)):,} sample)',
        transform=ax.transAxes, fontsize=9, color='#005B96', style='italic')
ax.set_xlim(-87.84, -87.52)
ax.set_ylim(41.64, 42.02)
ax.set_xlabel('Longitude')
ax.set_ylabel('Latitude')
ax.legend(loc='upper right', fontsize=9)
ax.set_aspect('equal')
plt.tight_layout(rect=[0, 0, 1, 0.96])
plt.savefig(f'{OUTPUT_DIR}/09_crash_density_map.png', bbox_inches='tight')
plt.close()
print("Saved: 09_crash_density_map.png")

# Chart 10: Year-over-Year Trend
yearly = df.groupby('CRASH_YEAR').agg(
    total_crashes=('HAS_INJURY', 'count'),
    injury_crashes=('HAS_INJURY', 'sum'),
    fatalities=('INJURIES_FATAL', 'sum')
).reset_index()

fig, ax = plt.subplots(figsize=(10, 5))
fig.patch.set_linewidth(0)
ax.plot(yearly['CRASH_YEAR'], yearly['total_crashes'], 'o-', color='#005B96',
        linewidth=2.5, markersize=8, markerfacecolor='white', markeredgewidth=2, label='Total Crashes')
ax.plot(yearly['CRASH_YEAR'], yearly['injury_crashes'], 's-', color='#E63946',
        linewidth=2, markersize=7, markerfacecolor='white', markeredgewidth=2, label='Injury Crashes')
ax.yaxis.grid(True, alpha=0.3, linestyle='--')
ax.yaxis.set_major_formatter(ticker.FuncFormatter(lambda x, _: f'{x:,.0f}'))
# COVID shaded band
ax.axvspan(2019.7, 2020.3, alpha=0.1, color='red', zorder=0)
covid_idx = yearly[yearly['CRASH_YEAR'] == 2020].index[0]
ax.annotate('COVID-19\nLockdown', xy=(2020, yearly.loc[covid_idx, 'total_crashes']),
           xytext=(2020.8, yearly['total_crashes'].max() * 0.85),
           fontsize=10, fontweight='bold', color='red',
           arrowprops=dict(arrowstyle='->', color='red', lw=1.5))
# Data value labels at each marker
for _, row in yearly.iterrows():
    ax.text(row['CRASH_YEAR'], row['total_crashes'] + yearly['total_crashes'].max() * 0.02,
            f"{int(row['total_crashes']):,}", ha='center', fontsize=8, color='#005B96', fontweight='bold')
# Year-over-year % change annotations
for i in range(1, len(yearly)):
    prev_val = yearly.iloc[i-1]['total_crashes']
    curr_val = yearly.iloc[i]['total_crashes']
    pct_change = (curr_val - prev_val) / prev_val * 100
    mid_y = (prev_val + curr_val) / 2
    mid_x = yearly.iloc[i-1]['CRASH_YEAR'] + 0.5
    change_color = '#E63946' if pct_change < 0 else '#2A9D8F'
    ax.text(mid_x, mid_y, f'{pct_change:+.0f}%', fontsize=7, color=change_color,
            fontweight='bold', ha='center', va='center',
            bbox=dict(boxstyle='round,pad=0.15', facecolor='white', alpha=0.8, edgecolor='none'))
ax.set_xlabel('Year')
ax.set_ylabel('Number of Crashes')
ax.set_title('Year-over-Year Crash Trend (2016\u20132023)', fontsize=14, fontweight='bold')
fig.text(0.5, 0.93, '2020 crash reduction reflects COVID-19 lockdowns; recovery visible from 2021',
         ha='center', fontsize=10, color='gray', style='italic')
# Annotate 2016 — check if it's actually a partial year
earliest_2016 = df[df['CRASH_YEAR']==2016]['CRASH_DATE'].min()
print(f"  Earliest 2016 crash date: {earliest_2016}")
if earliest_2016.month >= 6:
    ann_2016_text = '2016: partial year\n(data collection began mid-year)'
else:
    ann_2016_text = '2016: lower volume\n(early dataset period)'
ax.annotate(ann_2016_text,
           xy=(2016, yearly[yearly['CRASH_YEAR']==2016]['total_crashes'].values[0]),
           xytext=(2016.5, yearly[yearly['CRASH_YEAR']==2016]['total_crashes'].values[0] * 0.75),
           fontsize=8, color='gray', style='italic',
           arrowprops=dict(arrowstyle='->', color='gray', lw=1))
ax.legend(loc='upper left')
ax.set_xticks(yearly['CRASH_YEAR'])
plt.tight_layout(rect=[0, 0, 1, 0.92])
plt.savefig(f'{OUTPUT_DIR}/10_yearly_trend.png', bbox_inches='tight')
plt.close()
print("Saved: 10_yearly_trend.png")

# Interactive Maps
print("\nGenerating interactive maps...")
from folium.plugins import MiniMap, Fullscreen, MousePosition

# --- Map 1: Crash Heatmap with Hotspot Clusters ---
m1 = folium.Map(location=[41.88, -87.63], zoom_start=11, tiles='CartoDB positron')
Fullscreen().add_to(m1)
MiniMap(toggle_display=True).add_to(m1)
MousePosition(position='bottomleft', prefix='Coords:').add_to(m1)

# Title overlay
title_html = (
    '<div style="position:fixed;top:10px;left:60px;z-index:9999;'
    'background:rgba(0,91,150,0.92);padding:12px 24px;border-radius:8px;'
    'box-shadow:0 2px 8px rgba(0,0,0,0.3);">'
    '<span style="color:white;font-size:18px;font-weight:bold;font-family:Calibri,sans-serif;">'
    'Chicago Crash Density Heatmap with Hotspot Clusters</span><br>'
    '<span style="color:#c8dcf0;font-size:12px;font-family:Calibri,sans-serif;">'
    f'{total_crashes//1000}K crashes (2016&ndash;2023) &bull; 50K sample shown &bull; Red circles = DBSCAN injury hotspots</span>'
    '</div>'
)
m1.get_root().html.add_child(folium.Element(title_html))

heat_sample = valid_coords.sample(n=min(50000, len(valid_coords)), random_state=42)
heat_data = heat_sample[['LATITUDE', 'LONGITUDE']].values.tolist()

# Heatmap layer in its own feature group
fg_heat = folium.FeatureGroup(name='Crash Density Heatmap', show=True)
HeatMap(heat_data, radius=8, blur=10, max_zoom=13, gradient={0.2:'blue',0.4:'lime',0.6:'yellow',0.8:'orange',1.0:'red'}).add_to(fg_heat)
fg_heat.add_to(m1)

# Hotspot clusters in their own feature group
if len(hotspots_df) > 0:
    fg_hotspots = folium.FeatureGroup(name='Injury Hotspot Clusters (DBSCAN)', show=True)
    for idx, (_, row) in enumerate(hotspots_df.head(10).iterrows()):
        radius_px = max(8, row['crash_count'] / hotspots_df['crash_count'].max() * 25)
        popup_html = (
            f"<div style='font-family:Calibri,sans-serif;min-width:200px;'>"
            f"<b style='font-size:14px;color:#005B96;'>Hotspot #{idx+1}: {row['top_street']}</b><br><hr>"
            f"<b>Crashes:</b> {int(row['crash_count']):,}<br>"
            f"<b>Injuries:</b> {int(row['total_injuries']):,}<br>"
            f"<b>Fatalities:</b> {int(row['fatalities'])}<br>"
            f"<b>Top cause:</b> {row['top_cause']}<br>"
            f"<b>Avg speed limit:</b> {row['mean_speed_limit']:.0f} mph<br>"
            f"<b>Lighting:</b> {row['top_lighting']}</div>"
        )
        folium.CircleMarker(
            location=[row['center_lat'], row['center_lon']],
            radius=radius_px,
            color='red', fill=True, fill_color='red', fill_opacity=0.7, weight=2,
            popup=folium.Popup(popup_html, max_width=280),
            tooltip=f"#{idx+1} {row['top_street']} ({int(row['crash_count']):,} crashes)"
        ).add_to(fg_hotspots)
        # Number label
        folium.Marker(
            location=[row['center_lat'], row['center_lon']],
            icon=folium.DivIcon(html=f'<div style="font-size:11px;font-weight:bold;color:darkred;text-shadow:1px 1px white;">{idx+1}</div>',
                                icon_size=(20, 20), icon_anchor=(10, 10))
        ).add_to(fg_hotspots)
    fg_hotspots.add_to(m1)

# Legend overlay
legend_html = '''
<div style="position:fixed;bottom:30px;right:10px;z-index:9999;
     background:white;padding:12px 16px;border-radius:6px;border:1px solid #ccc;
     box-shadow:0 2px 6px rgba(0,0,0,0.2);font-family:Calibri,sans-serif;font-size:12px;">
  <b style="font-size:13px;">Legend</b><br>
  <span style="display:inline-block;width:60px;height:10px;
        background:linear-gradient(to right,blue,lime,yellow,orange,red);
        border-radius:2px;vertical-align:middle;margin-right:5px;"></span>
  Crash density (low &rarr; high)<br>
  <span style="display:inline-block;width:14px;height:14px;border-radius:50%;
        background:red;border:2px solid darkred;vertical-align:middle;margin-right:5px;"></span>
  Injury hotspot cluster<br>
  <span style="color:gray;font-size:10px;margin-top:4px;display:block;">
    Click any red circle for full details</span>
</div>
'''
m1.get_root().html.add_child(folium.Element(legend_html))
folium.LayerControl(collapsed=False).add_to(m1)
m1.save(f'{OUTPUT_DIR}/crash_heatmap.html')
print("Saved: crash_heatmap.html")

# --- Map 2: Fatal Crashes ---
fatal_crashes = df[(df['HAS_FATAL'] == 1) & df['LATITUDE'].notna() & df['LONGITUDE'].notna()]
m2 = folium.Map(location=[41.88, -87.63], zoom_start=11, tiles='CartoDB positron')
Fullscreen().add_to(m2)
MiniMap(toggle_display=True).add_to(m2)
MousePosition(position='bottomleft', prefix='Coords:').add_to(m2)

# Title overlay
title2_html = f'''
<div style="position:fixed;top:10px;left:60px;z-index:9999;
     background:rgba(139,0,0,0.92);padding:12px 24px;border-radius:8px;
     box-shadow:0 2px 8px rgba(0,0,0,0.3);">
  <span style="color:white;font-size:18px;font-weight:bold;font-family:Calibri,sans-serif;">
    Fatal Crash Locations &mdash; Chicago (2016&ndash;2023)</span><br>
  <span style="color:#f0c0c0;font-size:12px;font-family:Calibri,sans-serif;">
    {len(fatal_crashes):,} fatal crashes plotted &bull; Click any marker for details</span>
</div>
'''
m2.get_root().html.add_child(folium.Element(title2_html))

# Fatal crashes by year for layer grouping
fg_all = folium.FeatureGroup(name=f'All Fatal Crashes ({len(fatal_crashes):,})', show=True)
for _, row in fatal_crashes.iterrows():
    crash_yr = int(row['CRASH_DATE'].year) if hasattr(row['CRASH_DATE'], 'year') else 'N/A'
    cause = row['PRIM_CONTRIBUTORY_CAUSE'] if pd.notna(row['PRIM_CONTRIBUTORY_CAUSE']) else 'Unknown'
    speed = f"{int(row['POSTED_SPEED_LIMIT'])} mph" if pd.notna(row['POSTED_SPEED_LIMIT']) else 'N/A'
    popup_html = (
        f"<div style='font-family:Calibri,sans-serif;min-width:220px;'>"
        f"<b style='font-size:13px;color:darkred;'>Fatal Crash</b><br><hr>"
        f"<b>Date:</b> {str(row['CRASH_DATE'])[:10]}<br>"
        f"<b>Cause:</b> {cause[:50]}<br>"
        f"<b>Speed limit:</b> {speed}<br>"
        f"<b>Lighting:</b> {row['LIGHTING_CONDITION'] if pd.notna(row['LIGHTING_CONDITION']) else 'N/A'}<br>"
        f"<b>Weather:</b> {row['WEATHER_CONDITION'] if pd.notna(row['WEATHER_CONDITION']) else 'N/A'}</div>"
    )
    folium.CircleMarker(
        location=[row['LATITUDE'], row['LONGITUDE']],
        radius=5, color='darkred', fill=True, fill_color='darkred', fill_opacity=0.7, weight=1,
        popup=folium.Popup(popup_html, max_width=300),
        tooltip=f"Fatal: {str(row['CRASH_DATE'])[:10]}"
    ).add_to(fg_all)
fg_all.add_to(m2)

# Legend for fatal map
legend2_html = f'''
<div style="position:fixed;bottom:30px;right:10px;z-index:9999;
     background:white;padding:12px 16px;border-radius:6px;border:1px solid #ccc;
     box-shadow:0 2px 6px rgba(0,0,0,0.2);font-family:Calibri,sans-serif;font-size:12px;">
  <b style="font-size:13px;">Legend</b><br>
  <span style="display:inline-block;width:12px;height:12px;border-radius:50%;
        background:darkred;border:1px solid #600;vertical-align:middle;margin-right:5px;"></span>
  Fatal crash location<br>
  <span style="color:gray;font-size:11px;margin-top:4px;display:block;">
    Total: {len(fatal_crashes):,} fatal crashes (2016&ndash;2023)<br>
    Click any marker for date, cause &amp; conditions</span>
</div>
'''
m2.get_root().html.add_child(folium.Element(legend2_html))
folium.LayerControl(collapsed=False).add_to(m2)
m2.save(f'{OUTPUT_DIR}/fatal_crashes_map.html')
print(f"Saved: fatal_crashes_map.html ({len(fatal_crashes):,} fatal crashes)")

# --- Chart 11: Static Fatal Crash Map (for PPTX embedding) ---
fig, ax = plt.subplots(figsize=(10, 12))
fig.patch.set_linewidth(0)
# Background: all crashes as light density
ax.scatter(sample['LONGITUDE'], sample['LATITUDE'], alpha=0.015, s=0.5, c='#b0c4de', zorder=1)
# Fatal crashes as prominent dots
ax.scatter(fatal_crashes['LONGITUDE'], fatal_crashes['LATITUDE'],
           s=12, c='darkred', alpha=0.6, edgecolors='red', linewidth=0.3, zorder=3, label=f'Fatal crashes ({len(fatal_crashes):,})')
# Hotspot overlays
if len(hotspots_df) > 0:
    for idx, (_, row) in enumerate(hotspots_df.head(10).iterrows()):
        label = 'Injury hotspot centres (DBSCAN)' if idx == 0 else None
        ax.plot(row['center_lon'], row['center_lat'], 'o', color='none',
                markeredgecolor='orange', markersize=15, markeredgewidth=2, zorder=4,
                label=label)
# Geographic references
for lon, lat, name in [(-87.63, 41.882, 'Downtown'), (-87.63, 41.76, 'South Side'),
                        (-87.75, 41.98, "O'Hare Area"), (-87.72, 41.88, 'West Side')]:
    ax.text(lon, lat, name, fontsize=10, color='navy', alpha=0.4, fontweight='bold', ha='center', style='italic')
ax.set_xlim(-87.84, -87.52)
ax.set_ylim(41.64, 42.02)
ax.set_xlabel('Longitude')
ax.set_ylabel('Latitude')
ax.set_title('Fatal Crash Locations \u2014 Chicago (2016\u20132023)', fontsize=14, fontweight='bold')
fig.text(0.5, 0.97, f'{len(fatal_crashes):,} fatal crashes plotted. Orange rings = injury hotspot centres.',
         ha='center', fontsize=10, color='gray', style='italic')
ax.legend(loc='upper right', fontsize=9)
ax.set_aspect('equal')
plt.tight_layout(rect=[0, 0, 1, 0.96])
plt.savefig(f'{OUTPUT_DIR}/11_fatal_crash_map.png', bbox_inches='tight')
plt.close()
print("Saved: 11_fatal_crash_map.png")

# ============================================================
# PHASE 5: POWERPOINT PRESENTATION
# ============================================================
print("\n" + "=" * 70)
print("PHASE 5: POWERPOINT GENERATION")
print("=" * 70)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BLUE = RGBColor(0, 91, 150)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
RED = RGBColor(230, 57, 70)
LIGHT_GRAY = RGBColor(120, 120, 120)

def add_title_bar(slide, title_text, subtitle_text=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Calibri'
    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(12), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(200, 220, 240)
        p2.font.name = 'Calibri'

def add_body_text(slide, text, left=0.5, top=1.5, width=12, height=5.5, font_size=16):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
    return txBox

def add_image(slide, img_path, left=0.5, top=1.5, width=None, height=None):
    if os.path.exists(img_path):
        if width and height:
            slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))
        elif width:
            slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(width))
        elif height:
            slide.shapes.add_picture(img_path, Inches(left), Inches(top), height=Inches(height))

def add_section_header(question_text, answer_subtitle):
    """Dark blue full-bleed question-header slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = question_text
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = ""
    p2.font.size = Pt(14)
    p3 = tf.add_paragraph()
    p3.text = answer_subtitle
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(200, 220, 240)
    p3.font.name = 'Calibri'
    p3.alignment = PP_ALIGN.CENTER
    p3.font.italic = True
    return slide

# ================================================================
# SLIDE 1: Title
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()
txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Chicago Traffic Crash Analysis"
p.font.size = Pt(44)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "2016\u20132023"
p2.font.size = Pt(36)
p2.font.color.rgb = RGBColor(200, 220, 240)
p2.font.name = 'Calibri'
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = ""
p3.font.size = Pt(10)
p4 = tf.add_paragraph()
p4.text = "Data-Driven Insights for Road Safety"
p4.font.size = Pt(22)
p4.font.color.rgb = RGBColor(200, 220, 240)
p4.font.name = 'Calibri'
p4.alignment = PP_ALIGN.CENTER
p5 = tf.add_paragraph()
p5.text = ""
p5.font.size = Pt(20)
p6 = tf.add_paragraph()
p6.text = "Ali Maraci  |  February 2026"
p6.font.size = Pt(16)
p6.font.color.rgb = RGBColor(180, 200, 220)
p6.font.name = 'Calibri'
p6.alignment = PP_ALIGN.CENTER

# ================================================================
# SLIDE 2: Executive Summary
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Executive Summary",
              f"{total_crashes:,} crashes across Chicago (2016\u20132023) \u2014 {total_injuries:,} injuries, {total_fatalities:,} fatalities")

peak_hr_hour_val = hr_rate_hour.idxmax()
peak_hr_pct_val = hr_rate_hour.max()
peak_hr_label = f'{int(peak_hr_hour_val) % 12 or 12} {"AM" if peak_hr_hour_val < 12 else "PM"}'
summary = (
    f"This analysis answers four research questions using {total_crashes//1000}K crash records from the\n"
    "City of Chicago Open Data Portal, with the goal of identifying actionable patterns\n"
    "to inform road safety interventions.\n\n"
    f"Q1 \u2014 When do crashes happen?\n"
    f"    Crashes are not random. {calendar.month_name[peak_month]} is the worst month, "
    f"{peak_day_name} the worst day, and {int(peak_hour):02d}:00 the worst hour.\n"
    f"    Monthly, daily, and hourly patterns are all statistically significant (p < 0.001).\n"
    f"    Year-over-year trend is not significant (r={r_year:.2f}, p={p_year:.2f}).\n\n"
    f"Q2 \u2014 Is wet weather associated with more injuries?\n"
    f"    Wet weather raises the chance of injury by about {(wet_or - 1)*100:.0f}% "
    f"(adjusted Odds Ratio, OR), even after accounting for darkness, speed, time of day, and month.\n"
    f"    The unadjusted OR is higher at {OR:.2f} (+{(OR-1)*100:.0f}%). "
    f"This is a strong association, though we cannot prove direct causation from this data alone.\n\n"
    f"Q3 \u2014 What else stands out in the data?\n"
    f"    Late-night hit-and-run rates are elevated ({peak_hr_label} peak: {peak_hr_pct_val:.0f}%). "
    f"Higher speed limits dramatically increase fatality risk ({multiplier:.1f}x from lowest to highest zone). "
    f"'Failing to yield' is the single most dangerous driver behaviour.\n\n"
    f"Q4 \u2014 Where are the worst locations?\n"
    f"    {n_clusters} injury crash clusters identified using DBSCAN (Density-Based Spatial Clustering\n"
    f"    of Applications with Noise). {top_hotspot_street} corridor dominates.\n"
    f"    Targeted interventions recommended per cluster."
)
add_body_text(slide, summary, font_size=13)

# ================================================================
# SLIDE 3: Methodology & Tools
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Methodology & Tools",
              "Reproducible pipeline: raw CSV (Comma-Separated Values) \u2192 cleaned data \u2192 statistical tests \u2192 visualisations \u2192 this report")
methods = (
    "WHY THIS APPROACH:\n"
    "Traffic crash data is rich but noisy. We use a multi-method pipeline to ensure\n"
    "findings are robust: SQL for exploration, statistical tests for confirmation,\n"
    "regression for confounding control, and spatial clustering for location analysis.\n\n"
    "HOW WE ANALYSED THE DATA:\n\n"
    "PYTHON LIBRARIES: pandas, numpy, scipy, statsmodels, scikit-learn, matplotlib, seaborn, folium, python-pptx\n\n"
    "\u2022 Data cleaning & preparation  \u2014  pandas / numpy (filtering, missing values, feature creation)\n"
    "\u2022 SQL (Structured Query Language) exploration  \u2014  SQLite queries to identify initial trends\n"
    "\u2022 Statistical testing  \u2014  chi-squared, Mann-Whitney U, Spearman correlation\n"
    "   (These tests answer: 'Is this pattern real or just chance?')\n"
    "\u2022 Association analysis  \u2014  logistic regression controlling for confounders\n"
    "   (Isolates the effect of weather from darkness, speed, time-of-day, and month)\n"
    "\u2022 Hotspot detection  \u2014  DBSCAN (Density-Based Spatial Clustering of Applications with Noise)\n"
    "   geographic clustering (groups nearby crashes into meaningful clusters)\n"
    "\u2022 Mapping  \u2014  folium interactive maps + matplotlib static charts\n"
    "\u2022 Reporting  \u2014  python-pptx automated slide generation\n\n"
    "STATISTICAL METHODS USED:\n"
    "Chi-squared goodness-of-fit  |  Seasonal decomposition  |  Mann-Whitney U\n"
    "Odds Ratio (OR) with 95% Confidence Interval (CI)  |  Logistic regression  |  Stratified analysis\n"
    "DBSCAN clustering  |  Point-biserial correlation  |  Spearman rank correlation\n\n"
    f"DATA: {total_crashes:,} crash records, {len(df.columns)} variables, City of Chicago 2016\u20132023"
)
add_body_text(slide, methods, font_size=13)

# ================================================================
# SLIDE 4: Data Overview
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Data Overview",
              "City of Chicago Open Data Portal \u2014 Traffic Crashes dataset")
overview = (
    "DATA SOURCE:\n"
    "The City of Chicago Open Data Portal publishes crash-level records from police reports.\n"
    "Each row represents one crash event with details on timing, location, severity, and conditions.\n\n"
    "WHAT'S IN THE DATA:\n"
    f"\u2022 {total_crashes:,} individual crash records from 2016 to 2023\n"
    f"\u2022 {len(df.columns)} columns: when it happened, where, how severe, weather, road conditions, cause\n"
    "\u2022 Geographic scope: entire City of Chicago\n\n"
    "KEY NUMBERS AT A GLANCE:\n"
    f"\u2022 1 in {injury_ratio} crashes results in an injury (injury rate: {injury_rate:.1f}%)\n"
    f"\u2022 Fatality rate: {df['HAS_FATAL'].mean()*100:.3f}% ({total_fatalities:,} deaths over 8 years)\n"
    f"\u2022 Hit-and-run: {hr_pct:.1f}% of all crashes ({int(len(hr)):,} incidents)\n"
    f"\u2022 Wet weather: {df['IS_WET_WEATHER'].mean()*100:.1f}% of crashes occur during rain, snow, or freezing conditions\n\n"
    "DATA QUALITY DECISIONS:\n"
    "\u2022 Dropped DOORING_I (99.7% blank) and LANE_CNT (74.7% missing) \u2014 too incomplete for analysis\n"
    "\u2022 540K blank hit-and-run values conservatively treated as 'No'\n"
    "\u2022 Created 12 derived features (wet weather flag, severity flag, time period, etc.)"
)
add_body_text(slide, overview, font_size=13)

# ================================================================
# SLIDE 5: Q1 Section Header
# ================================================================
add_section_header(
    "Q1: Is there a correlation between time of year, week, and day and crashes?",
    "Yes \u2014 strong, statistically significant temporal patterns exist across every time dimension"
)

# ================================================================
# SLIDE 6: Q1 Monthly Seasonality
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Q1: Monthly Seasonality & Seasonal Decomposition",
              f"Crashes follow a clear annual rhythm \u2014 {peak_month_name} peak, {trough_month_name} low")
add_image(slide, f'{OUTPUT_DIR}/01_time_series_decomposition.png', left=0.3, top=1.4, width=8)
add_body_text(slide,
    "WHY THIS MATTERS:\n"
    "Understanding seasonal patterns\n"
    "allows agencies to allocate\n"
    "enforcement and resources to\n"
    "match when crashes peak.\n\n"
    "WHAT THIS SHOWS:\n"
    "Additive seasonal decomposition\n"
    "splits monthly crash counts into\n"
    "four layers:\n"
    "  1. Observed (raw data)\n"
    "  2. Long-term trend\n"
    "  3. Repeating seasonal cycle\n"
    "  4. Residual (random variation)\n\n"
    "KEY TAKEAWAYS:\n"
    f"\u2022 {peak_month_name} is the worst month\n"
    f"  (\u03c7\u00b2 = {chi2_month:,.0f}, p < 0.001)\n\n"
    "\u2022 Trend panel clearly shows the\n"
    "  COVID-19 dip and recovery\n\n"
    "\u2022 Seasonal cycle confirms\n"
    "  resources should be allocated\n"
    "  to match annual rhythm",
    left=8.5, top=1.5, width=4.5, font_size=12)

# ================================================================
# SLIDE 7: Q1 Day & Hour Patterns
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Q1: Day & Hour Crash Patterns",
              f"{peak_day_name} afternoon rush hour is the single most dangerous time window")
add_image(slide, f'{OUTPUT_DIR}/02_dow_hour_heatmap.png', left=0.3, top=1.4, width=8.5)
add_body_text(slide,
    "WHY THIS MATTERS:\n"
    "Knowing exactly when crashes\n"
    "peak allows targeted patrol\n"
    "scheduling and public messaging.\n\n"
    "WHAT THIS SHOWS:\n"
    "A heatmap where darker cells\n"
    "mean more crashes. Each row is\n"
    "a day; each column is an hour.\n\n"
    "KEY TAKEAWAYS:\n"
    f"\u2022 {peak_day_name} {peak_hour_label_main} = highest risk\n"
    "\u2022 Weekday rush hours dominate\n"
    "\u2022 Saturday nights show a\n"
    "  late-night activity spike\n"
    f"\u2022 {safest_label} is safest\n\n"
    "STATISTICAL CONFIRMATION:\n"
    f"\u2022 Hourly \u03c7\u00b2 = {chi2_hour:,.0f}\n"
    f"\u2022 Day-of-week \u03c7\u00b2 = {chi2_dow:,.0f}\n"
    "  (Both p < 0.001)",
    left=9, top=1.5, width=4, font_size=12)

# ================================================================
# SLIDE 8: Q1 Year-over-Year Trend
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Q1: Year-over-Year Crash Trend",
              "COVID-19 drove a sharp 2020 drop; crashes have since rebounded")
add_image(slide, f'{OUTPUT_DIR}/10_yearly_trend.png', left=0.3, top=1.4, width=8.5)
add_body_text(slide,
    "WHY THIS MATTERS:\n"
    "Long-term trends help distinguish\n"
    "structural changes from one-off\n"
    "events (e.g. the COVID-19 shock).\n\n"
    "WHAT THIS SHOWS:\n"
    "Annual total and injury crash\n"
    "counts from 2016 to 2023, with\n"
    "year-on-year % change.\n\n"
    "KEY TAKEAWAYS:\n"
    "\u2022 2020 saw a sharp crash drop\n"
    "  due to COVID-19 lockdowns\n\n"
    "\u2022 Post-2020: steady recovery\n"
    "  towards pre-pandemic levels\n\n"
    "\u2022 2016 is low (partial/early\n"
    "  dataset period)\n\n"
    "\u2022 Injury crashes track the\n"
    f"  same pattern at ~{injury_rate:.0f}% of total",
    left=9, top=1.5, width=4, font_size=12)

# ================================================================
# SLIDE 9: Q2 Section Header
# ================================================================
add_section_header(
    "Q2: Is there an association between wet weather conditions and crash severity?",
    "Strong association confirmed through 6 levels of evidence; true causation requires further study"
)

# ================================================================
# SLIDE 10: Q2 Weather Distribution
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Q2: Crash Volume & Injury Rate by Weather",
              "Clear weather has the most crashes (it's the most common); rain and freezing have higher injury rates")
add_image(slide, f'{OUTPUT_DIR}/04_weather_impact.png', left=0.3, top=1.4, width=8.5)
add_body_text(slide,
    "WHY THIS MATTERS:\n"
    "If wet weather is linked to more\n"
    "injuries, it justifies investment\n"
    "in drainage, anti-skid surfaces,\n"
    "and dynamic warning signs.\n\n"
    "WHAT THIS SHOWS:\n"
    "Horizontal bars = number of\n"
    "crashes per weather type.\n"
    "Red text = % of those crashes\n"
    "that resulted in an injury.\n\n"
    "KEY TAKEAWAYS:\n"
    "\u2022 Most crashes happen in clear\n"
    "  weather (it's the most common)\n\n"
    "\u2022 But rain, freezing rain and\n"
    "  fog have HIGHER injury rates\n"
    f"  ({rain_inj_rate:.0f}%+ vs {clear_inj_rate:.0f}% for clear)\n\n"
    "\u2022 Next slide adjusts for how\n"
    "  often each weather occurs",
    left=9, top=1.5, width=4, font_size=11)

# ================================================================
# SLIDE 11: Q2 Exposure-Adjusted
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Q2: Exposure-Adjusted Analysis",
              "Comparing crash share to weather share reveals over/under-representation")
exposure = (
    "THE PROBLEM WITH RAW COUNTS:\n"
    "Clear weather has the most crashes simply because most days are clear.\n"
    "To be fair, we compare average daily crash counts on wet vs. dry days.\n\n"
    "EXPOSURE-ADJUSTED COMPARISON:\n"
    "Using crash-date weather records (NULL/UNKNOWN excluded):\n\n"
    f"                     Days      Avg Crashes/Day    CRI (Crash Rate Index)\n"
    f"  Wet Weather     {n_wet_days:>5,}       {avg_crashes_wet_day:>8.1f}            {cri_wet:.2f}\n"
    f"  Dry Weather     {n_dry_days:>5,}       {avg_crashes_dry_day:>8.1f}            1.00 (baseline)\n\n"
    "WHAT THIS MEANS:\n"
    f"\u2022 CRI (Crash Rate Index) = {cri_wet:.2f}: wet weather days have "
    f"{'fewer' if cri_wet < 1 else 'more'} crashes per day than dry days\n"
    "\u2022 But when crashes DO occur in wet weather, they're more likely to cause injuries\n"
    "\u2022 This distinction between crash frequency and crash severity is critical\n\n"
    "IMPORTANT CAVEAT:\n"
    "'Wet days' are classified by whether any crash that day had a wet-weather code, not by\n"
    "independent weather station data (e.g. NOAA (National Oceanic and Atmospheric\n"
    "Administration) O'Hare). Many crashes on a 'wet day' likely occurred during dry hours.\n"
    "The CRI therefore reflects reporting patterns, not a true weather-adjusted crash rate.\n"
    "Self-selection bias also applies: cautious drivers may stay home in bad weather."
)
add_body_text(slide, exposure, font_size=13)

# ================================================================
# SLIDE 12: Q2 Logistic Regression
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Q2: Isolating Weather's Effect \u2014 Logistic Regression",
              "After controlling for darkness, speed, time, and month: wet weather still raises injury risk")
add_image(slide, f'{OUTPUT_DIR}/05_odds_ratio_chart.png', left=0.3, top=1.4, width=7)
wet_or_row = or_df[or_df['Variable'] == 'IS_WET_WEATHER'].iloc[0]
dark_or_row = or_df[or_df['Variable'] == 'IS_DARK'].iloc[0]
spd_or_row = or_df[or_df['Variable'] == 'POSTED_SPEED_LIMIT'].iloc[0]
add_body_text(slide,
    "WHY LOGISTIC REGRESSION:\n"
    "Simple comparisons can be\n"
    "misleading (wet crashes may\n"
    "also happen in darkness or at\n"
    "higher speeds). Regression\n"
    "isolates each factor's\n"
    "independent contribution.\n\n"
    "WHAT THIS SHOWS:\n"
    "An Odds Ratio (OR) chart. An\n"
    "OR > 1 means higher injury risk.\n\n"
    "HOW TO READ IT:\n"
    f"\u2022 Wet weather: OR = {wet_or_row['OR']:.3f}\n"
    f"  \u2192 ~{(wet_or_row['OR']-1)*100:.0f}% higher injury odds\n\n"
    f"\u2022 Darkness: OR = {dark_or_row['OR']:.3f}\n"
    f"  \u2192 ~{(dark_or_row['OR']-1)*100:.0f}% higher injury odds\n\n"
    f"\u2022 Speed limit (+1 mph): OR = {spd_or_row['OR']:.3f}\n"
    f"  \u2192 Each extra mph adds ~{(spd_or_row['OR']-1)*100:.0f}%\n\n"
    f"All p < 0.001 | n = {nobs:,}\n"
    f"Pseudo R\u00b2 = {r2:.3f} (low R\u00b2 is\n"
    f"typical for crash-level models;\n"
    f"individual ORs remain reliable)",
    left=7.5, top=1.4, width=5.5, font_size=11)

# ================================================================
# SLIDE 13: Q2 Conclusions & Limitations
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Q2: Weather Conclusions & Limitations",
              "Six independent lines of evidence all point the same way")
add_body_text(slide,
    "WHAT WE FOUND (6 LEVELS OF EVIDENCE):\n\n"
    f"\u2713 1. Naive comparison \u2014 Wet weather crashes have higher injury rates (+{(wet_inj_rate - dry_inj_rate)*100:.1f} pp)\n"
    "\u2713 2. Exposure-adjusted \u2014 Fewer crashes per wet day, but injuries are worse when they happen\n"
    "\u2713 3. Chi-squared test \u2014 The injury rate difference is statistically significant (not chance)\n"
    f"\u2713 4. Unadjusted Odds Ratio (OR) \u2014 OR = {OR:.2f} (+{(OR-1)*100:.0f}% higher injury odds)\n"
    f"\u2713 5. Adjusted logistic regression \u2014 Adjusted OR = {wet_or:.3f} (+{(wet_or-1)*100:.0f}%), controlling for darkness, speed, time, and month\n"
    "\u2713 6. Stratified check \u2014 The pattern holds whether it's daylight, dark, or dusk\n\n"
    "BOTTOM LINE:\n"
    "Wet weather is consistently linked to more severe crashes. The evidence is strong,\n"
    "but this is an OBSERVATIONAL STUDY \u2014 we can demonstrate association, not prove causation.\n\n"
    "ON THE QUESTION OF CAUSATION:\n"
    "The task asks specifically about a causal relationship. Proving causation from observational\n"
    "data requires quasi-experimental methods (e.g., instrumental variables, difference-in-differences,\n"
    "regression discontinuity). We recommend a further study matching daily weather station data\n"
    "(NOAA O'Hare) to crash rates, or before/after analysis of road surface treatments.\n\n"
    "WHAT WE CAN'T ACCOUNT FOR:\n"
    "\u2022 We don't know daily traffic volumes (fewer people may drive in bad weather)\n"
    "\u2022 Weather data is from crash reports, not independent weather station observations\n"
    "\u2022 Cautious drivers may stay home in storms (self-selection bias)\n"
    "\u2022 Visibility, road salt, and tire condition are unmeasured",
    font_size=13)

# ================================================================
# SLIDE 14: Q3 Section Header
# ================================================================
add_section_header(
    "Q3: Explore 2\u20133 key findings of your choice from the data",
    "Hit-and-run patterns  |  Speed\u2013severity relationship  |  Most dangerous driver behaviours"
)

# ================================================================
# SLIDE 15: Q3 Finding 1 — Hit-and-Run
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Q3 Finding 1: Hit-and-Run Crash Patterns",
              "Late-night crashes have elevated hit-and-run rates \u2014 possibly linked to impairment or reduced witnesses")
add_image(slide, f'{OUTPUT_DIR}/06_hit_and_run.png', left=0.3, top=1.4, width=8.5)
add_body_text(slide,
    "WHY WE INVESTIGATED THIS:\n"
    f"~{hr_pct:.0f}% of crashes are hit-and-runs\n"
    "\u2014 far higher than most cities.\n"
    "Understanding the temporal pattern\n"
    "can guide enforcement strategy.\n\n"
    "WHAT THIS SHOWS:\n"
    "Left: Hit-and-run rate by hour.\n"
    "Right: Hit-and-run rate by day.\n"
    "Dashed line = overall average.\n\n"
    "KEY NUMBERS:\n"
    f"\u2022 {len(hr):,} hit-and-runs\n"
    f"  ({len(hr)/len(df)*100:.0f}% of all crashes)\n\n"
    f"\u2022 {peak_hr_label} peak: {peak_hr_pct_val:.0f}% of crashes\n"
    "  are hit-and-run\n\n"
    "SO WHAT?\n"
    "\u2022 Late-night enforcement should\n"
    "  be increased (midnight\u20134 AM)\n"
    "\u2022 CCTV (Closed-Circuit Television)\n"
    "  cameras on high hit-and-run\n"
    "  corridors would help solve cases\n\n"
    "DATA NOTE:\n"
    "Hit-and-run flag is 'Y' or blank.\n"
    "540K blanks treated as 'No'.",
    left=9, top=1.5, width=4, font_size=11)

# ================================================================
# SLIDE 16: Q3 Finding 2 — Speed & Severity
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Q3 Finding 2: Speed Limit & Crash Severity",
              f"Higher speed limits are associated with {multiplier:.0f}x higher fatality rates")
add_image(slide, f'{OUTPUT_DIR}/07_speed_severity.png', left=0.3, top=1.4, width=8.5)
fat_min = speed_stats['fatality_rate_per_1000'].iloc[0]
fat_max = speed_stats['fatality_rate_per_1000'].iloc[-1]
multiplier = fat_max / fat_min if fat_min > 0 else 0
add_body_text(slide,
    "WHY THIS MATTERS:\n"
    "Speed management is one of the\n"
    "most effective levers for reducing\n"
    "crash fatalities. This data shows\n"
    "a clear dose-response pattern.\n\n"
    "WHAT THIS SHOWS:\n"
    "Red bars = injury rate (left axis).\n"
    "Blue line = fatality rate per 1,000\n"
    "crashes (right axis).\n\n"
    "KEY NUMBERS:\n"
    f"\u2022 Fatality rate is {multiplier:.1f}x higher\n"
    "  in 45+ mph zones vs 0\u201320 mph\n\n"
    f"\u2022 Spearman \u03c1 = {rho:.3f} (p < 0.001)\n"
    "  = speed and severity are linked\n\n"
    "SO WHAT?\n"
    "\u2022 Speed cameras in 35+ mph zones\n"
    "\u2022 Consider road diets on high-\n"
    "  crash corridors (reduce lanes)\n"
    "\u2022 Lower limits near schools/parks",
    left=9, top=1.5, width=4, font_size=11)

# ================================================================
# SLIDE 17: Q3 Finding 3 — Contributory Causes
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Q3 Finding 3: Most Dangerous Driver Behaviours",
              "Ranked by total injury crashes \u2014 identifies behaviours that are both common and harmful")
add_image(slide, f'{OUTPUT_DIR}/08_contributory_causes.png', left=0.3, top=1.4, width=8.5)
add_body_text(slide,
    "WHY THIS MATTERS:\n"
    "Identifying which driver behaviours\n"
    "cause the most harm focuses\n"
    "enforcement and education on\n"
    "the highest-impact targets.\n\n"
    "WHAT THIS SHOWS:\n"
    "The 15 most dangerous causes,\n"
    "ranked by total injury crashes.\n"
    "Colour = injury rate\n"
    "(red = more dangerous).\n\n"
    "TOP 3 MOST DANGEROUS:\n"
    "\u2022 Failing to yield right-of-way\n"
    "\u2022 Following too closely\n"
    "\u2022 Failing to reduce speed\n\n"
    "SO WHAT?\n"
    "\u2022 Protected turn phases at\n"
    "  high-yield-failure intersections\n"
    "\u2022 Public awareness campaigns\n"
    "\u2022 Enforcement + red-light cameras",
    left=9, top=1.5, width=4, font_size=11)

# ================================================================
# SLIDE 18: Q4 Section Header
# ================================================================
add_section_header(
    "Q4: Where do most crashes occur and what safety measures should be considered?",
    f"{n_clusters} injury crash clusters identified via spatial analysis \u2014 targeted interventions for each"
)

# ================================================================
# SLIDE 19: Q4 Hotspot Map (static)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Q4: Crash Hotspot Map \u2014 Injury Clusters",
              "DBSCAN clustering groups nearby injury crashes to reveal dangerous corridors")
add_image(slide, f'{OUTPUT_DIR}/09_crash_density_map.png', left=0.3, top=1.3, width=6, height=5.8)
hotspot_text = (
    "WHY SPATIAL CLUSTERING:\n"
    "Street names alone don't pinpoint\n"
    "WHERE on a street the problems\n"
    "are. DBSCAN (Density-Based Spatial\n"
    "Clustering of Applications with\n"
    "Noise) automatically identifies\n"
    "geographic clusters of nearby\n"
    "crashes without pre-specifying\n"
    "the number of clusters.\n\n"
    "HOW IT WORKS:\n"
    f"Groups crashes within {eps_m}m of\n"
    f"each other (min {min_samp} crashes).\n"
    f"Result: {n_clusters} distinct clusters.\n\n"
    "TOP 10 INJURY HOTSPOTS:\n"
)
for i, row in hotspots_df.head(10).iterrows():
    hotspot_text += f"{i+1}. {row['top_street']} \u2014 {int(row['crash_count']):,} crashes\n"
hotspot_text += (
    "\n\u27a4 Interactive version available:\n"
    "   crash_heatmap.html\n"
    "   (open in any web browser for\n"
    "   full zoom, click-for-details)"
)
add_body_text(slide, hotspot_text, left=6.5, top=1.3, width=6.5, font_size=12)

# ================================================================
# SLIDE 20: Q4 Fatal Crash Map (static)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Q4: Fatal Crash Locations",
              f"{len(fatal_crashes):,} fatal crashes mapped ({total_fatalities:,} deaths) \u2014 concentrated on high-speed arterials and the South/West sides")
add_image(slide, f'{OUTPUT_DIR}/11_fatal_crash_map.png', left=0.3, top=1.3, width=6, height=5.8)
add_body_text(slide,
    "WHAT THIS SHOWS:\n"
    "Every fatal crash from 2016\u20132023\n"
    "plotted on a map of Chicago.\n"
    "Orange rings = injury hotspot centres.\n\n"
    "KEY PATTERNS:\n"
    "\u2022 Fatal crashes cluster on major\n"
    f"  arterials ({top_4_streets})\n\n"
    "\u2022 South and West sides are\n"
    "  disproportionately affected\n\n"
    "\u2022 Higher-speed corridors show\n"
    "  more fatal incidents\n\n"
    "\u27a4 Interactive version available:\n"
    "   fatal_crashes_map.html\n"
    "   (click any dot for date, cause,\n"
    "   speed limit, weather, lighting)",
    left=6.5, top=1.3, width=6.5, font_size=12)

# ================================================================
# SLIDE 21: Q4 Safety Recommendations
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Q4: Recommended Safety Interventions",
              "Each finding maps directly to an actionable countermeasure")
add_body_text(slide,
    "FINDING \u2192 ACTION:\n\n"
    f"Speed kills (45+ mph zones have {multiplier:.1f}x fatality rate):\n"
    "   \u2192 Reduce speed limits on identified corridors, deploy speed cameras,\n"
    "      implement road diets (4-lane to 3-lane with turn lane)\n\n"
    f"Darkness worsens outcomes (+{(dark_or_row['OR']-1)*100:.0f}% injury odds):\n"
    "   \u2192 LED (Light-Emitting Diode) street lighting upgrades at the highest-injury\n"
    "      intersections, high-visibility crosswalks with flashing beacons\n\n"
    "'Failing to yield' is the #1 dangerous behaviour:\n"
    "   \u2192 Protected turn phases, roundabout pilot studies, red-light cameras\n"
    "      at intersections with the highest yield-failure crash counts\n\n"
    f"Hit-and-runs peak midnight\u20134 AM ({peak_hr_pct_val:.0f}% at {peak_hr_label}):\n"
    "   \u2192 CCTV (Closed-Circuit Television) on high hit-and-run corridors,\n"
    "      late-night enforcement, public awareness campaigns\n\n"
    f"Wet weather increases injury risk by ~{(wet_or-1)*100:.0f}%:\n"
    "   \u2192 Improved drainage, anti-skid surface treatments, dynamic warning signs\n\n"
    f"{peak_month_name} is peak crash month:\n"
    f"   \u2192 Seasonal safety campaigns in {calendar.month_name[peak_month - 1] if peak_month > 1 else calendar.month_name[12]}/{peak_month_name}, enforcement surges",
    font_size=13)

# ================================================================
# SLIDE 22: Further Steps
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Further Steps & Recommendations",
              "What additional data and analyses would strengthen these findings")
add_body_text(slide,
    "TO MAKE THE ANALYSIS MORE RIGOROUS:\n\n"
    "1. TRAFFIC VOLUME DATA  \u2014  Obtain AADT (Annual Average Daily Traffic) counts from IDOT\n"
    "   (Illinois Department of Transportation) to calculate true crash rates per vehicle-mile\n\n"
    "2. DAILY WEATHER RECORDS  \u2014  Match each crash to actual weather station data from NOAA\n"
    "   (National Oceanic and Atmospheric Administration) for precise exposure adjustment\n\n"
    "3. VEHICLE & PERSON DATA  \u2014  Link to companion datasets for driver age,\n"
    "   impairment, vehicle type, seatbelt usage\n\n"
    "TO TAKE ACTION:\n\n"
    "4. BEFORE/AFTER STUDIES  \u2014  Measure the impact of past interventions\n"
    "   (speed cameras, road redesigns) using interrupted time series analysis\n\n"
    "5. PREDICTIVE MODEL  \u2014  Random Forest or XGBoost to score crash risk\n"
    "   by location + time + weather for proactive resource deployment\n\n"
    "6. VULNERABLE ROAD USERS  \u2014  Separate analysis for pedestrians and cyclists\n\n"
    "7. ECONOMIC COSTS  \u2014  Attach FHWA (Federal Highway Administration) cost estimates\n"
    "   ($1.7M per fatality, $98K per injury) to prioritise interventions by\n"
    "   ROI (Return on Investment)",
    font_size=13)

# ================================================================
# SLIDE 23: Appendix — SQL Queries
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Appendix: SQL Queries Used")
add_body_text(slide,
    "-- Crashes by Year\n"
    "SELECT CRASH_YEAR, COUNT(*) as total_crashes,\n"
    "       SUM(INJURIES_TOTAL) as total_injuries,\n"
    "       ROUND(AVG(HAS_INJURY)*100, 1) as injury_rate_pct\n"
    "FROM crashes GROUP BY CRASH_YEAR ORDER BY CRASH_YEAR;\n\n"
    "-- Weather Impact\n"
    "SELECT WEATHER_CONDITION, COUNT(*) as total_crashes,\n"
    "       ROUND(AVG(HAS_INJURY)*100, 1) as injury_rate_pct\n"
    "FROM crashes WHERE WEATHER_CONDITION != 'UNKNOWN'\n"
    "GROUP BY WEATHER_CONDITION ORDER BY total_crashes DESC;\n\n"
    "-- Top Contributory Causes\n"
    "SELECT PRIM_CONTRIBUTORY_CAUSE, COUNT(*) as total_crashes,\n"
    "       ROUND(AVG(HAS_INJURY)*100, 1) as injury_rate_pct\n"
    "FROM crashes WHERE PRIM_CONTRIBUTORY_CAUSE NOT IN\n"
    "('UNABLE TO DETERMINE', 'NOT APPLICABLE')\n"
    "GROUP BY PRIM_CONTRIBUTORY_CAUSE ORDER BY total_crashes DESC LIMIT 15;\n\n"
    "-- Day x Hour Cross-tab | Hit-and-Run Overview | Top 20 Streets",
    font_size=11)

# ================================================================
# SLIDE 24: Appendix — Interactive Maps
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Appendix: Interactive Maps (HTML files)",
              "Two zoomable, clickable maps are included as separate HTML (HyperText Markup Language) files")
add_body_text(slide,
    "These maps cannot be embedded directly in PowerPoint.\n"
    "Open the HTML files in any web browser (Chrome, Edge, Firefox) for full interactivity.\n\n"
    "1. crash_heatmap.html  \u2014  CRASH DENSITY HEATMAP\n"
    "   \u2022 50,000 crash locations shown as a colour-coded density layer\n"
    "   \u2022 Red circles mark the top 10 injury hotspot clusters\n"
    "   \u2022 Click any red circle for: crash count, injuries, fatalities, top cause, speed limit\n"
    "   \u2022 Gradient legend, layer toggle, minimap, and fullscreen mode included\n"
    "   \u2022 Zoom in to explore individual neighbourhoods\n\n"
    "2. fatal_crashes_map.html  \u2014  FATAL CRASH LOCATIONS\n"
    f"   \u2022 All {len(fatal_crashes):,} fatal crashes plotted as dark-red markers\n"
    "   \u2022 Click any marker for: date, cause, speed limit, lighting, weather conditions\n"
    "   \u2022 Enables identification of fatal crash corridors and clusters\n\n"
    "Both maps use CartoDB Positron tiles for a clean, readable base map.\n"
    "The static versions of these maps are shown on Slides 19 and 20.",
    font_size=14)

# Save PowerPoint
output_path = f'{OUTPUT_DIR}/Chicago_Traffic_Crash_Analysis.pptx'
prs.save(output_path)
print(f"\nPowerPoint saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")

# Close SQLite
conn.close()
print("SQLite connection closed.")

# Summary
print("\n" + "=" * 70)
print("ANALYSIS COMPLETE")
print("=" * 70)
print("\nOutput files:")
for f in sorted(os.listdir(OUTPUT_DIR)):
    fpath = os.path.join(OUTPUT_DIR, f)
    size = os.path.getsize(fpath)
    print(f"  {f} ({size/1024:.0f} KB)")
